"""
Semantic Alignment Engine
Maps Whisper transcript segments → slide pages using dense RAG.

Pipeline:
  1. Extract text from each slide (PDF / PPTX / DOCX)
     - Includes speaker notes (PPTX) and image descriptions (OCR + Claude vision)
  2. Embed slides with sentence-transformers (all-mpnet-base-v2, GPU)
  3. Build a FAISS index for fast cosine-similarity lookup
  4. For each transcript segment query the index (optionally with a context
     window of ±CONTEXT_SEC seconds for richer matching signal)
  5. Apply Viterbi temporal smoothing so slides only advance forward
  6. Flag segments that don't match any slide well (off_slide)
  7. Collapse consecutive equal-slide segments into a compact timeline
  8. Save JSON to [course_id]/alignment/[stem].json

Usage:
  # align one caption↔slide pair
  python semantic_alignment.py \\
      --caption  85427/captions/CS3210\\ e-Lecture\\ on\\ Processes\\ and\\ Threads.json \\
      --slides   85427/materials/LectureNotes/L02-Processes-Threads.pdf

  # auto-discover all unaligned pairs in a course folder
  python semantic_alignment.py --course 85427
"""

from __future__ import annotations

import argparse
import io
import json
import re
import sys
from collections import defaultdict
from pathlib import Path
from typing import NamedTuple

try:
    import faiss
    import numpy as np
except ImportError as _e:
    print(f"[error] Missing dependency: {_e}")
    print("[error] Please install the ML environment from Settings → ML Environment in the AutoNote app.")
    sys.exit(1)

PROJECT_DIR = Path(__file__).parent
import sys as _sys
_AUTO_NOTE_DIR = Path.home() / ".auto_note"
import os as _os
if _os.environ.get("AUTONOTE_DATA_DIR"):
    DATA_DIR = Path(_os.environ["AUTONOTE_DATA_DIR"])
elif getattr(_sys, "frozen", False) or PROJECT_DIR == _AUTO_NOTE_DIR / "scripts":
    DATA_DIR = _AUTO_NOTE_DIR
else:
    DATA_DIR = PROJECT_DIR

# Course output directory: defaults to DATA_DIR but can be overridden by
# OUTPUT_DIR in config.json so files land in the user's chosen Output Dir.
try:
    _cfg_file = DATA_DIR / "config.json"
    _sa_config: dict = json.loads(_cfg_file.read_text(encoding="utf-8")) if _cfg_file.exists() else {}
except Exception:
    _sa_config = {}
_out_dir = _sa_config.get("OUTPUT_DIR", "").strip()
COURSE_DATA_DIR = Path(_out_dir) if _out_dir else DATA_DIR

# ── Tunable knobs ─────────────────────────────────────────────────────────────

EMBED_MODEL   = "all-mpnet-base-v2"   # highest-quality general sentence model
CONTEXT_SEC   = 30.0                  # seconds of transcript to pool per query
BATCH_SIZE    = 64                    # embedding batch size

# Jina Embeddings v4 multimodal model — used when aligning transcript to slide
# images directly (no slide text extraction needed).  Requires the jina API key
# or the local model jinaai/jina-embeddings-v4 via sentence-transformers.
JINA_EMBED_MODEL = "jinaai/jina-embeddings-v4"
JINA_API_URL     = "https://api.jina.ai/v1/embeddings"

# Viterbi transition log-probabilities
STAY_LOGP     =  0.0    # free to stay on the same slide
FWD_LOGP_PER  = -0.02   # cost per slide advanced forward (nearly free)
BWD_LOGP_PER  = -1.5    # cost per slide stepped backward
                         # firm — allows genuine multi-minute review sections
                         # (e.g. slide 63 revisited) while suppressing
                         # single-segment noise flips

# Temporal position prior: at time t, add a Gaussian bonus centred on the
# expected slide position (t / duration) * n_slides.
# PRIOR_SIGMA controls the width in slide units; larger = softer guide.
PRIOR_SIGMA   = 8.0

# Off-slide detection: segments where the best raw cosine similarity (before
# the position prior) is below this threshold are flagged as off_slide.
OFF_SLIDE_THRESHOLD = 0.28

# Image captioning: pages / slides with fewer than this many words of text
# have their visual content described via OCR + Claude vision.
IMAGE_WORD_THRESHOLD = 20

# OpenAI vision model used for slide image description.
# gpt-4o-mini is used for cost efficiency; each call is only made once per
# slide because results are cached in {slide_file}.image_cache.json.
OPENAI_VISION_MODEL = "gpt-4o-mini"

# Sparse-slide enrichment: slides below this word count borrow text from
# content-rich neighbours before embedding.
SPARSE_THRESHOLD = 30   # words
NEIGHBOR_WORDS   = 60   # how many words to borrow from each neighbour


# ── OpenAI API key helper ─────────────────────────────────────────────────────

def _get_openai_key() -> str:
    import os
    key = os.environ.get("OPENAI_API_KEY", "")
    if not key:
        key_file = DATA_DIR / "openai_api.txt"
        if key_file.exists():
            key = key_file.read_text().strip()
    return key


# ── Image description (OCR + Claude vision) ───────────────────────────────────

_VISION_PROMPT = (
    "This is a slide from a university lecture. "
    "Describe ALL visible content in detail: every text label, concept name, "
    "diagram element, arrow, relationship, numbered step, code snippet, and "
    "technical term you can see. "
    "Focus on the academic topic — mention state names, algorithm steps, "
    "data structure relationships, or any specific terminology shown. "
    "Write a thorough plain-text description (no markdown), 3-5 sentences."
)


class ImageDescriber:
    """
    Extracts textual descriptions from slide images using:
      1. pytesseract OCR  — free, captures printed text in diagrams
      2. Claude vision API — powerful, understands academic diagrams and
                             technical content (used only on slides with
                             < IMAGE_WORD_THRESHOLD words; results cached)

    The Claude API is called at most once per slide per slide file, thanks to
    the image cache written to {slide_file}.image_cache.json.
    """

    def __init__(self):
        self._ocr_ok    = None   # None = unchecked, True/False after first call
        self._openai    = None   # openai.OpenAI client, lazy-loaded
        self._api_calls = 0      # count for cost reporting

    # ── lazy loaders ─────────────────────────────────────────────────────────

    def _ensure_ocr(self) -> bool:
        if self._ocr_ok is None:
            try:
                import pytesseract
                pytesseract.get_tesseract_version()
                self._ocr_ok = True
            except Exception:
                print("  [img] pytesseract not available — OCR skipped")
                self._ocr_ok = False
        return self._ocr_ok

    def _ensure_openai(self) -> bool:
        if self._openai is not None:
            return True
        key = _get_openai_key()
        if not key:
            print("  [img] No OpenAI API key — vision description skipped")
            return False
        try:
            from openai import OpenAI
            self._openai = OpenAI(api_key=key)
            return True
        except Exception as e:
            print(f"  [img] OpenAI client error: {e}")
            return False

    # ── per-image methods ────────────────────────────────────────────────────

    def _ocr(self, img) -> str:
        if not self._ensure_ocr():
            return ""
        try:
            import pytesseract
            return pytesseract.image_to_string(img).strip()
        except Exception:
            return ""

    def _openai_describe(self, img) -> str:
        """Send one slide image to GPT-4o-mini vision and return description."""
        if not self._ensure_openai():
            return ""
        import base64
        try:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.standard_b64encode(buf.getvalue()).decode()

            resp = self._openai.chat.completions.create(
                model=OPENAI_VISION_MODEL,
                max_tokens=300,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url",
                         "image_url": {"url": f"data:image/png;base64,{b64}",
                                       "detail": "low"}},   # "low" = cheapest
                        {"type": "text", "text": _VISION_PROMPT},
                    ],
                }],
            )
            self._api_calls += 1
            return resp.choices[0].message.content.strip()
        except Exception as e:
            print(f"  [img] OpenAI vision error: {e}")
            return ""

    # ── public API ───────────────────────────────────────────────────────────

    def describe_slide_image(self, img) -> str:
        """
        Run OCR then Claude vision on a single slide image (already rendered
        as a full-page pixmap).  Returns combined description text.
        The caller is responsible for caching so this is never called twice
        for the same slide.
        """
        from PIL import Image as PILImage
        if not isinstance(img, PILImage.Image):
            return ""
        parts: list[str] = []
        ocr = self._ocr(img)
        if ocr:
            parts.append(ocr)
        vision = self._openai_describe(img)
        if vision:
            parts.append(vision)
        return " ".join(parts)

    @property
    def api_calls(self) -> int:
        return self._api_calls


# ── Image cache helpers ────────────────────────────────────────────────────────

def _load_image_cache(slide_path: Path) -> dict:
    cache_file = slide_path.parent / f"{slide_path.name}.image_cache.json"
    if cache_file.exists():
        with open(cache_file) as f:
            return json.load(f)
    return {}


def _save_image_cache(slide_path: Path, cache: dict) -> None:
    cache_file = slide_path.parent / f"{slide_path.name}.image_cache.json"
    with open(cache_file, "w") as f:
        json.dump(cache, f, indent=2)


# ── Slide text extraction ─────────────────────────────────────────────────────

class SlideText(NamedTuple):
    index: int          # 0-based slide/page index
    label: str          # short title (first non-empty line)
    text:  str          # full extracted text (incl. notes + image descriptions)


def extract_pdf(path: Path,
                describer: ImageDescriber | None = None,
                img_cache: dict | None = None) -> list[SlideText]:
    import fitz
    from PIL import Image as PILImage

    doc    = fitz.open(str(path))
    slides = []
    for i, page in enumerate(doc):
        text  = page.get_text().strip()
        label = next((ln.strip() for ln in text.splitlines() if ln.strip()),
                     f"Page {i+1}")

        # Image enrichment for sparse pages — render the full page as a
        # high-res pixmap so vector diagrams (state graphs, flowcharts, etc.)
        # are captured alongside any embedded bitmaps.
        if describer is not None and len(text.split()) < IMAGE_WORD_THRESHOLD:
            cache_key = f"page_{i}"
            if img_cache is not None and cache_key in img_cache:
                img_desc = img_cache[cache_key]
            else:
                pxmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                pil   = PILImage.frombytes(
                    "RGB", [pxmap.width, pxmap.height], pxmap.samples)
                img_desc = describer.describe_slide_image(pil)
                if img_cache is not None:
                    img_cache[cache_key] = img_desc

            if img_desc:
                text = (text + "\n" + img_desc).strip()

        slides.append(SlideText(i, label[:80], text))
    doc.close()
    return slides


def extract_pptx(path: Path,
                 describer: ImageDescriber | None = None,
                 img_cache: dict | None = None) -> list[SlideText]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image as PILImage

    prs    = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []

        # Shape text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

        # Speaker notes
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(notes_text)
        except Exception:
            pass

        text  = "\n".join(parts)
        label = parts[0][:80] if parts else f"Slide {i+1}"

        # Image enrichment for sparse slides: use the first picture shape.
        # For PPTX we extract the embedded image blob directly (no rendering).
        if describer is not None and len(text.split()) < IMAGE_WORD_THRESHOLD:
            cache_key = f"slide_{i}"
            if img_cache is not None and cache_key in img_cache:
                img_desc = img_cache[cache_key]
            else:
                img_desc = ""
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            pil = PILImage.open(
                                io.BytesIO(shape.image.blob)).convert("RGB")
                            img_desc = describer.describe_slide_image(pil)
                            break   # one image per slide is enough
                        except Exception:
                            pass
                if img_cache is not None:
                    img_cache[cache_key] = img_desc

            if img_desc:
                text = (text + "\n" + img_desc).strip()

        slides.append(SlideText(i, label, text))
    return slides


def extract_docx(path: Path) -> list[SlideText]:
    """
    Word documents lack explicit page breaks; we treat each paragraph as one
    logical unit and group them into synthetic 'pages' of ≈ PAGE_PARA paragraphs.
    """
    from docx import Document
    PAGE_PARA = 15
    doc   = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages: list[SlideText] = []
    for pi, start in enumerate(range(0, max(len(paras), 1), PAGE_PARA)):
        chunk = paras[start:start + PAGE_PARA]
        text  = "\n".join(chunk)
        label = chunk[0][:80] if chunk else f"Page {pi+1}"
        pages.append(SlideText(pi, label, text))
    return pages


def load_slides(path: Path,
                describer: ImageDescriber | None = None,
                img_cache: dict | None = None) -> list[SlideText]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path, describer, img_cache)
    if ext in (".pptx", ".ppt"):
        return extract_pptx(path, describer, img_cache)
    if ext in (".docx", ".doc"):
        return extract_docx(path)
    raise ValueError(f"Unsupported slide format: {ext}")


# ── Embedding & FAISS index ───────────────────────────────────────────────────

def get_embedder():
    """Load sentence-transformer model once; use GPU if available."""
    from sentence_transformers import SentenceTransformer
    import torch
    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"  [embed] Loading {EMBED_MODEL} on {device} ...", flush=True)
    try:
        model = SentenceTransformer(EMBED_MODEL, device=device)
    except Exception as e:
        print(f"  [embed] Failed to load model: {e}", flush=True)
        raise
    print(f"  [embed] Model loaded.", flush=True)
    return model


def embed_texts(model, texts: list[str], desc: str = "  embedding") -> np.ndarray:
    """Return L2-normalised float32 embeddings, shape (N, D)."""
    print(f"{desc} ({len(texts)} texts)...", flush=True)
    vecs = model.encode(
        texts,
        batch_size=BATCH_SIZE,
        show_progress_bar=False,     # avoid tqdm issues in Windows pipe-mode subprocesses
        normalize_embeddings=True,   # cosine sim → inner product on unit sphere
        convert_to_numpy=True,
    )
    print(f"{desc} done.", flush=True)
    return vecs.astype(np.float32)


def build_faiss_index(embeddings: np.ndarray) -> faiss.IndexFlatIP:
    """Inner-product index (= cosine similarity on normalised vectors)."""
    dim   = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    return index


# ── Jina Embeddings v4 multimodal ────────────────────────────────────────────

def _get_jina_key() -> str:
    import os
    key = os.environ.get("JINA_API_KEY", "")
    if not key:
        key_file = DATA_DIR / "jina_api.txt"
        if key_file.exists():
            key = key_file.read_text().strip()
    return key


def embed_images_jina(image_paths: list[Path], desc: str = "  embedding images") -> np.ndarray | None:
    """Embed slide images using Jina Embeddings v4 API.

    Returns L2-normalised float32 embeddings, shape (N, D), or None on failure.
    """
    import base64
    import requests as _req

    key = _get_jina_key()
    if not key:
        print("  [jina] No Jina API key — cannot embed images")
        return None

    print(f"{desc} ({len(image_paths)} images via Jina API)...", flush=True)

    # Encode images as base64
    inputs = []
    for p in image_paths:
        if not p.exists():
            continue
        b64 = base64.b64encode(p.read_bytes()).decode("ascii")
        inputs.append({"image": f"data:image/png;base64,{b64}"})

    if not inputs:
        return None

    # Call Jina API in batches
    all_embeddings = []
    batch_size = 16
    for i in range(0, len(inputs), batch_size):
        batch = inputs[i:i + batch_size]
        try:
            resp = _req.post(
                JINA_API_URL,
                headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
                json={"model": "jina-embeddings-v4", "input": batch,
                      "dimensions": 768, "normalized": True,
                      "embedding_type": "float", "task": "retrieval.passage"},
                timeout=120,
            )
            if resp.status_code != 200:
                print(f"  [jina] API error {resp.status_code}: {resp.text[:200]}")
                return None
            data = resp.json()
            for item in data.get("data", []):
                all_embeddings.append(item["embedding"])
        except Exception as e:
            print(f"  [jina] Request failed: {e}")
            return None

    print(f"{desc} done.", flush=True)
    return np.array(all_embeddings, dtype=np.float32)


def embed_texts_jina(texts: list[str], desc: str = "  embedding texts") -> np.ndarray | None:
    """Embed texts using Jina Embeddings v4 API.

    Returns L2-normalised float32 embeddings, shape (N, D), or None on failure.
    """
    import requests as _req

    key = _get_jina_key()
    if not key:
        print("  [jina] No Jina API key — cannot embed texts")
        return None

    print(f"{desc} ({len(texts)} texts via Jina API)...", flush=True)

    all_embeddings = []
    batch_size = 64
    for i in range(0, len(texts), batch_size):
        batch = [{"text": t} for t in texts[i:i + batch_size]]
        try:
            resp = _req.post(
                JINA_API_URL,
                headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
                json={"model": "jina-embeddings-v4", "input": batch,
                      "dimensions": 768, "normalized": True,
                      "embedding_type": "float", "task": "retrieval.query"},
                timeout=120,
            )
            if resp.status_code != 200:
                print(f"  [jina] API error {resp.status_code}: {resp.text[:200]}")
                return None
            data = resp.json()
            for item in data.get("data", []):
                all_embeddings.append(item["embedding"])
        except Exception as e:
            print(f"  [jina] Request failed: {e}")
            return None

    print(f"{desc} done.", flush=True)
    return np.array(all_embeddings, dtype=np.float32)


# ── Transcript windowing ──────────────────────────────────────────────────────

def build_window_texts(segments: list[dict], context_sec: float) -> list[str]:
    """
    For each segment build a richer query string by pooling the text of all
    segments whose midpoint falls within ±context_sec of this segment's midpoint.
    This prevents very short segments (e.g. "Yes" or "Okay") from misleading
    the matcher.
    """
    mids = [(s["start"] + s["end"]) / 2.0 for s in segments]
    n    = len(segments)
    texts: list[str] = []
    for i, mid in enumerate(mids):
        parts = []
        # scan backwards
        j = i
        while j >= 0 and mids[j] >= mid - context_sec:
            j -= 1
        # scan forwards
        k = i
        while k < n and mids[k] <= mid + context_sec:
            k += 1
        for s in segments[j+1 : k]:
            if s["text"].strip():
                parts.append(s["text"].strip())
        texts.append(" ".join(parts) if parts else (segments[i]["text"] or " "))
    return texts


# ── Viterbi temporal smoothing ────────────────────────────────────────────────

def viterbi_smooth(
    log_likelihoods: np.ndarray,   # shape (T, N_slides)
) -> list[int]:
    """
    Viterbi decoding that prefers forward progression.

    log_likelihoods[t, s]  = log P(observation_t | slide s)
                           = cosine_similarity score (already log-scale proxy)
    Returns: list of best slide indices, length T.
    """
    T, N = log_likelihoods.shape

    # dp[t, s] = best total log-score ending at slide s at step t
    dp   = np.full((T, N), -np.inf, dtype=np.float64)
    back = np.zeros((T, N), dtype=np.int32)

    dp[0] = log_likelihoods[0]

    for t in range(1, T):
        for s in range(N):
            # consider all previous states s_prev
            trans = np.zeros(N, dtype=np.float64)
            for sp in range(N):
                delta = s - sp
                if delta == 0:
                    trans[sp] = STAY_LOGP
                elif delta > 0:
                    trans[sp] = FWD_LOGP_PER * delta
                else:
                    trans[sp] = BWD_LOGP_PER * abs(delta)
            scores       = dp[t-1] + trans
            best_prev    = int(np.argmax(scores))
            dp[t, s]     = scores[best_prev] + log_likelihoods[t, s]
            back[t, s]   = best_prev

    # Traceback
    path     = [0] * T
    path[-1] = int(np.argmax(dp[-1]))
    for t in range(T - 2, -1, -1):
        path[t] = int(back[t + 1, path[t + 1]])
    return path


def viterbi_smooth_fast(log_likelihoods: np.ndarray) -> list[int]:
    """
    Vectorised Viterbi — O(T × N) instead of O(T × N²).
    Equivalent to the loop version but ~100× faster.
    """
    T, N = log_likelihoods.shape

    dp   = np.full((T, N), -np.inf, dtype=np.float64)
    back = np.zeros((T, N), dtype=np.int32)
    dp[0] = log_likelihoods[0]

    # Precompute transition matrix  trans[s_prev, s]
    idx   = np.arange(N)
    delta = idx[None, :] - idx[:, None]   # (N, N)  delta[sp, s] = s - sp
    trans = np.where(delta == 0, STAY_LOGP,
            np.where(delta >  0, FWD_LOGP_PER * delta,
                                 BWD_LOGP_PER * np.abs(delta)))

    for t in range(1, T):
        # scores[sp, s] = dp[t-1, sp] + trans[sp, s]
        scores        = dp[t-1, :, None] + trans          # (N, N)
        best_prev     = np.argmax(scores, axis=0)          # (N,)
        dp[t]         = scores[best_prev, np.arange(N)] + log_likelihoods[t]
        back[t]       = best_prev

    path     = [0] * T
    path[-1] = int(np.argmax(dp[-1]))
    for t in range(T - 2, -1, -1):
        path[t] = int(back[t + 1, path[t + 1]])
    return path


# ── Timeline collapse ─────────────────────────────────────────────────────────

def build_timeline(segments: list[dict], slide_path: list[int],
                   slides: list[SlideText],
                   off_slide_mask: list[bool] | None = None) -> list[dict]:
    """
    Merge consecutive segments assigned to the same slide into one interval.
    Off-slide segments are excluded from the timeline.
    Returns list of {slide_1based, start, end, label}.
    """
    if not segments:
        return []

    if off_slide_mask is None:
        off_slide_mask = [False] * len(segments)

    timeline: list[dict] = []
    cur_slide: int | None = None
    cur_start = 0.0
    cur_end   = 0.0

    for i in range(len(segments)):
        s  = segments[i]
        si = slide_path[i]

        if off_slide_mask[i]:
            # Flush current span before the gap
            if cur_slide is not None:
                timeline.append({
                    "slide":  cur_slide + 1,
                    "start":  round(cur_start, 3),
                    "end":    round(cur_end, 3),
                    "label":  slides[cur_slide].label,
                })
                cur_slide = None
            continue

        if cur_slide is None:
            cur_slide = si
            cur_start = s["start"]
            cur_end   = s["end"]
        elif si == cur_slide:
            cur_end = s["end"]
        else:
            timeline.append({
                "slide":  cur_slide + 1,
                "start":  round(cur_start, 3),
                "end":    round(cur_end, 3),
                "label":  slides[cur_slide].label,
            })
            cur_slide = si
            cur_start = s["start"]
            cur_end   = s["end"]

    if cur_slide is not None:
        timeline.append({
            "slide":  cur_slide + 1,
            "start":  round(cur_start, 3),
            "end":    round(cur_end, 3),
            "label":  slides[cur_slide].label,
        })
    return timeline


# ── Sparse-slide enrichment ───────────────────────────────────────────────────

def _enrich_sparse_slides(texts: list[str]) -> list[str]:
    """
    Slides that are nearly empty (section headers, diagram-only, code-only)
    get almost no embedding signal.  Enrich them by appending up to
    NEIGHBOR_WORDS from the nearest content-rich neighbours so the embedder
    has something to work with.  The original label text is kept at the front
    so the slide's own identity still dominates.
    """
    word_counts = [len(t.split()) for t in texts]
    enriched    = list(texts)

    for i, (text, wc) in enumerate(zip(texts, word_counts)):
        if wc >= SPARSE_THRESHOLD:
            continue
        extra: list[str] = []
        for delta in (1, -1, 2, -2, 3, -3):
            j = i + delta
            if 0 <= j < len(texts) and word_counts[j] >= SPARSE_THRESHOLD:
                words = texts[j].split()[:NEIGHBOR_WORDS]
                extra.extend(words)
            if len(extra) >= NEIGHBOR_WORDS * 2:
                break
        if extra:
            enriched[i] = text + " " + " ".join(extra)

    return enriched


# ── Main alignment routine ────────────────────────────────────────────────────

def align(caption_path: Path, slide_path: Path,
          out_dir: Path, embedder=None) -> Path:
    """
    Full alignment pipeline for one (caption, slide) pair.
    Returns path of the saved JSON.
    """
    print(f"\n{'='*70}")
    print(f"Caption : {caption_path.name}")
    print(f"Slides  : {slide_path.name}")
    print(f"{'='*70}")

    # ── Load inputs ──────────────────────────────────────────────────────────
    with open(caption_path, encoding="utf-8") as f:
        caption = json.load(f)
    segments: list[dict] = caption["segments"]
    if not segments:
        print("  [skip] Caption has no segments.")
        return out_dir

    print(f"  Transcript: {len(segments)} segments, {caption['duration']:.0f}s")

    # ── Load slides (with image captioning + speaker notes) ───────────────────
    describer = ImageDescriber()
    img_cache = _load_image_cache(slide_path)
    cache_size_before = len(img_cache)

    print("  Extracting slide text (+ images + speaker notes)...")
    slides = load_slides(slide_path, describer, img_cache)
    print(f"  Slides: {len(slides)} pages/slides")

    # Save image cache if any new descriptions were generated
    new_descriptions = len(img_cache) - cache_size_before
    if new_descriptions > 0:
        _save_image_cache(slide_path, img_cache)
        print(f"  [img] {describer.api_calls} GPT-4o-mini vision API call(s) made "
              f"({new_descriptions} new slide(s) described, cached for future runs)")

    # ── Embed slides ──────────────────────────────────────────────────────────
    if embedder is None:
        embedder = get_embedder()

    slide_texts = _enrich_sparse_slides([s.text if s.text.strip() else s.label for s in slides])
    print(f"  Embedding {len(slides)} slides...")
    slide_embs = embed_texts(embedder, slide_texts, desc="  slides")
    index      = build_faiss_index(slide_embs)

    # ── Embed transcript segments (with context window) ───────────────────────
    print(f"  Building context windows (±{CONTEXT_SEC:.0f}s)...")
    window_texts = build_window_texts(segments, CONTEXT_SEC)

    print(f"  Embedding {len(window_texts)} transcript windows...")
    seg_embs = embed_texts(embedder, window_texts, desc="  transcript")   # (T, D)

    # ── Raw similarity scores: each segment vs all slides ─────────────────────
    # faiss.search returns (scores, indices); with IndexFlatIP + normalised
    # vectors scores are cosine similarities in [-1, 1].
    print("  Querying FAISS index...")
    k          = len(slides)
    sims, idxs = index.search(seg_embs, k)   # both (T, N_slides)

    # Build log-likelihood matrix: reorder sims by slide index
    T          = len(segments)
    N          = len(slides)
    log_ll     = np.zeros((T, N), dtype=np.float64)
    for t in range(T):
        for rank in range(k):
            log_ll[t, idxs[t, rank]] = float(sims[t, rank])

    # Save raw log-likelihoods for off-slide detection (before the prior)
    log_ll_raw = log_ll.copy()

    # ── Temporal position prior ───────────────────────────────────────────────
    # At segment t, the professor is expected to be near slide
    # expected_s = (t_mid / total_duration) * (N - 1).
    # A Gaussian prior around this expected position gently pushes the
    # Viterbi toward the correct slide when raw similarities are ambiguous
    # (e.g., consecutive slides with nearly identical content).
    total_duration = caption["duration"] or 1.0
    slide_idx      = np.arange(N, dtype=np.float64)
    for t, seg in enumerate(segments):
        t_mid      = (seg["start"] + seg["end"]) / 2.0
        expected_s = (t_mid / total_duration) * (N - 1)
        prior      = -0.5 * ((slide_idx - expected_s) / PRIOR_SIGMA) ** 2
        log_ll[t] += prior

    # ── Off-slide detection ───────────────────────────────────────────────────
    # Segments where the best raw cosine is below threshold are flagged.
    # The professor may be speaking off-topic, answering questions, or doing
    # live demos that don't correspond to any slide.
    raw_max_sim    = log_ll_raw.max(axis=1)   # (T,)
    off_slide_mask = (raw_max_sim < OFF_SLIDE_THRESHOLD).tolist()
    n_off = sum(off_slide_mask)
    if n_off:
        print(f"  Off-slide segments: {n_off}/{T} "
              f"(raw cosine < {OFF_SLIDE_THRESHOLD})")

    # ── Viterbi ───────────────────────────────────────────────────────────────
    print("  Running Viterbi smoothing...")
    slide_path_idx = viterbi_smooth_fast(log_ll)

    # ── Per-segment results ───────────────────────────────────────────────────
    aligned_segments = []
    for i, seg in enumerate(segments):
        si        = slide_path_idx[i]
        off       = off_slide_mask[i]
        raw_sim   = round(float(log_ll_raw[i, si]), 4)
        aligned_segments.append({
            "id":          seg["id"],
            "start":       seg["start"],
            "end":         seg["end"],
            "text":        seg["text"],
            "slide":       None if off else si + 1,    # 1-based; null = off-slide
            "slide_label": None if off else slides[si].label,
            "similarity":  raw_sim,
            "off_slide":   off,
        })

    timeline = build_timeline(segments, slide_path_idx, slides, off_slide_mask)

    # ── Output ────────────────────────────────────────────────────────────────
    result = {
        "lecture":        caption_path.stem,
        "slide_file":     slide_path.name,
        "total_slides":   len(slides),
        "total_segments": len(segments),
        "off_slide_count": n_off,
        "duration":       caption["duration"],
        "language":       caption.get("language", ""),
        "embed_model":    EMBED_MODEL,
        "context_sec":    CONTEXT_SEC,
        "segments":       aligned_segments,
        "timeline":       timeline,
    }

    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / f"{caption_path.stem}.json"
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    # Summary
    slide_counts = {}
    for i, si in enumerate(slide_path_idx):
        if not off_slide_mask[i]:
            slide_counts[si] = slide_counts.get(si, 0) + 1
    covered = len(slide_counts)
    print(f"  Covered {covered}/{N} slides across {len(timeline)} intervals")
    print(f"  Saved → {out_file}")
    return out_file


# ── Multi-slide alignment ─────────────────────────────────────────────────────

def align_multi_slides(
    caption_path: Path,
    slide_paths: list[Path],
    out_dir: Path,
    embedder=None,
) -> list[Path]:
    """
    Align one caption against one or more slide files.

    For a single slide file this is identical to align().
    For multiple slide files:
      - Slides are concatenated into one global index (file 0 first, then file 1, …)
      - A single Viterbi pass runs over the combined index
      - Results are split back per file with local (1-based) slide numbering
      - One output JSON is written per slide file, named {slide_stem}.json
        so note_generation._find_alignment() matches each file perfectly

    Returns list of output JSON paths (one per slide file).
    """
    if len(slide_paths) == 1:
        out = align(caption_path, slide_paths[0], out_dir, embedder)
        return [out] if isinstance(out, Path) else []

    print(f"\n{'='*70}")
    print(f"Caption : {caption_path.name}")
    print(f"Slides  : {[p.name for p in slide_paths]}")
    print(f"{'='*70}")

    # ── Load caption ──────────────────────────────────────────────────────────
    with open(caption_path, encoding="utf-8") as f:
        caption = json.load(f)
    segments: list[dict] = caption["segments"]
    if not segments:
        print("  [skip] Caption has no segments.")
        return []
    print(f"  Transcript: {len(segments)} segments, {caption['duration']:.0f}s")

    # ── Load all slide files; track per-file offsets ──────────────────────────
    describer    = ImageDescriber()
    all_slides:  list[SlideText] = []
    file_offsets: list[int]      = []   # global start index for each file
    file_sizes:   list[int]      = []   # number of slides in each file
    img_caches:   list[dict]     = []

    for sp in slide_paths:
        ic = _load_image_cache(sp)
        ic_before = len(ic)
        file_slides = load_slides(sp, describer, ic)
        if len(ic) > ic_before:
            _save_image_cache(sp, ic)
        file_offsets.append(len(all_slides))
        file_sizes.append(len(file_slides))
        img_caches.append(ic)
        all_slides.extend(file_slides)
        print(f"  {sp.name}: {len(file_slides)} slides (global offset {file_offsets[-1]})")

    N = len(all_slides)
    print(f"  Combined: {N} slides total")

    if describer.api_calls:
        print(f"  [img] {describer.api_calls} GPT-4o-mini call(s) made")

    # ── Embed combined slides ─────────────────────────────────────────────────
    if embedder is None:
        embedder = get_embedder()

    slide_texts = _enrich_sparse_slides(
        [s.text if s.text.strip() else s.label for s in all_slides]
    )
    print(f"  Embedding {N} combined slides...")
    slide_embs = embed_texts(embedder, slide_texts)
    index      = build_faiss_index(slide_embs)

    # ── Embed transcript windows ──────────────────────────────────────────────
    print(f"  Building context windows (±{CONTEXT_SEC:.0f}s)...")
    window_texts = build_window_texts(segments, CONTEXT_SEC)
    print(f"  Embedding {len(window_texts)} transcript windows...")
    seg_embs = embed_texts(embedder, window_texts, desc="  transcript")

    # ── FAISS search → log-likelihood matrix ──────────────────────────────────
    print("  Querying FAISS index...")
    T          = len(segments)
    sims, idxs = index.search(seg_embs, N)
    log_ll     = np.zeros((T, N), dtype=np.float64)
    for t in range(T):
        for rank in range(N):
            log_ll[t, idxs[t, rank]] = float(sims[t, rank])
    log_ll_raw = log_ll.copy()

    # ── Temporal position prior ───────────────────────────────────────────────
    total_duration = caption["duration"] or 1.0
    slide_idx      = np.arange(N, dtype=np.float64)
    for t, seg in enumerate(segments):
        t_mid      = (seg["start"] + seg["end"]) / 2.0
        expected_s = (t_mid / total_duration) * (N - 1)
        prior      = -0.5 * ((slide_idx - expected_s) / PRIOR_SIGMA) ** 2
        log_ll[t] += prior

    # ── Off-slide detection ───────────────────────────────────────────────────
    raw_max_sim    = log_ll_raw.max(axis=1)
    off_slide_mask = (raw_max_sim < OFF_SLIDE_THRESHOLD).tolist()
    n_off          = sum(off_slide_mask)
    if n_off:
        print(f"  Off-slide segments: {n_off}/{T} (raw cosine < {OFF_SLIDE_THRESHOLD})")

    # ── Viterbi ───────────────────────────────────────────────────────────────
    print("  Running Viterbi smoothing...")
    global_path = viterbi_smooth_fast(log_ll)

    # ── Split results per file and save ───────────────────────────────────────
    out_dir.mkdir(parents=True, exist_ok=True)
    out_files: list[Path] = []

    for fi, (sp, off, size) in enumerate(zip(slide_paths, file_offsets, file_sizes)):
        file_slides = all_slides[off : off + size]

        aligned_segments = []
        for i, seg in enumerate(segments):
            g_si    = global_path[i]
            off_seg = off_slide_mask[i]
            raw_sim = round(float(log_ll_raw[i, g_si]), 4)

            # Segment belongs to this file if its global slide index is in range
            in_file = (off <= g_si < off + size) and not off_seg
            local_si = (g_si - off) if in_file else None

            aligned_segments.append({
                "id":          seg["id"],
                "start":       seg["start"],
                "end":         seg["end"],
                "text":        seg["text"],
                "slide":       (local_si + 1) if in_file else None,
                "slide_label": file_slides[local_si].label if in_file else None,
                "similarity":  raw_sim,
                "off_slide":   not in_file,
            })

        # Build per-file Viterbi path (clamp to local range for timeline)
        local_path = [
            max(0, min(global_path[i] - off, size - 1))
            for i in range(T)
        ]
        # off_slide mask for timeline: original off_slide OR assigned to other file
        local_off = [
            off_slide_mask[i] or not (off <= global_path[i] < off + size)
            for i in range(T)
        ]
        timeline = build_timeline(segments, local_path, file_slides, local_off)

        covered = len({s["slide"] for s in aligned_segments if s["slide"] is not None})
        print(f"  {sp.name}: {covered}/{size} slides covered, "
              f"{len(timeline)} timeline intervals")

        result = {
            "lecture":         caption_path.stem,
            "slide_file":      sp.name,
            "slide_files_all": [p.name for p in slide_paths],
            "file_idx":        fi + 1,
            "total_slides":    size,
            "total_segments":  T,
            "off_slide_count": sum(1 for s in aligned_segments if s["off_slide"]),
            "duration":        caption["duration"],
            "language":        caption.get("language", ""),
            "embed_model":     EMBED_MODEL,
            "context_sec":     CONTEXT_SEC,
            "segments":        aligned_segments,
            "timeline":        timeline,
        }

        # Name by slide stem so _find_alignment() in note_generation.py scores 1.0
        out_file = out_dir / f"{sp.stem}.json"
        with open(out_file, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"  Saved → {out_file}")
        out_files.append(out_file)

    return out_files


# ── Jina v4 multimodal alignment (transcript text ↔ slide images) ────────────

def align_multimodal(
    caption_path: Path,
    slide_path: Path,
    out_dir: Path,
) -> Path | None:
    """Align transcript to slide images using Jina Embeddings v4.

    Instead of embedding slide text, this renders each slide page to an image
    and embeds it with Jina v4.  Transcript segments are embedded as text.
    Both live in the same multimodal vector space for cross-modal matching.

    Temporal heuristic rewards:
      - Consecutive segments get a bonus for staying on the same slide
      - Forward progression is preferred over backward jumps
      - Temporal position prior guides the Viterbi decoder
    """
    print(f"\n{'='*70}")
    print(f"Multimodal alignment (Jina v4)")
    print(f"Caption : {caption_path.name}")
    print(f"Slides  : {slide_path.name}")
    print(f"{'='*70}")

    # ── Load caption ─────────────────────────────────────────────────────────
    with open(caption_path, encoding="utf-8") as f:
        caption = json.load(f)
    segments: list[dict] = caption["segments"]
    if not segments:
        print("  [skip] Caption has no segments.")
        return None
    print(f"  Transcript: {len(segments)} segments, {caption['duration']:.0f}s")

    # ── Render slide images ──────────────────────────────────────────────────
    import tempfile
    tmp_img_dir = Path(tempfile.mkdtemp(prefix="jina_slides_"))
    print(f"  Rendering slides to images...")

    slide_images: list[Path] = []
    ext = slide_path.suffix.lower()
    if ext == ".pdf":
        import fitz
        from PIL import Image as PILImage
        doc = fitz.open(str(slide_path))
        for i in range(len(doc)):
            png = tmp_img_dir / f"slide_{i+1:03d}.png"
            px = doc[i].get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            pil = PILImage.frombytes("RGB", [px.width, px.height], px.samples)
            pil.save(str(png))
            slide_images.append(png)
            del px, pil
        doc.close()
    else:
        # For non-PDF, fall back to text-based alignment
        print("  [warn] Jina multimodal requires PDF slides; falling back to text alignment")
        return None

    N = len(slide_images)
    print(f"  Rendered {N} slide images")

    # ── Embed slide images with Jina v4 ──────────────────────────────────────
    slide_embs = embed_images_jina(slide_images, desc="  slides (Jina v4)")
    if slide_embs is None:
        print("  [error] Failed to embed slide images; falling back to text alignment")
        import shutil
        shutil.rmtree(tmp_img_dir, ignore_errors=True)
        return None

    # ── Embed transcript segments as text with Jina v4 ───────────────────────
    window_texts = build_window_texts(segments, CONTEXT_SEC)
    seg_embs = embed_texts_jina(window_texts, desc="  transcript (Jina v4)")
    if seg_embs is None:
        print("  [error] Failed to embed transcript; falling back to text alignment")
        import shutil
        shutil.rmtree(tmp_img_dir, ignore_errors=True)
        return None

    # ── Build similarity matrix ──────────────────────────────────────────────
    T = len(segments)
    # Cosine similarity (both are L2-normalised)
    log_ll = (seg_embs @ slide_embs.T).astype(np.float64)  # (T, N)
    log_ll_raw = log_ll.copy()

    # ── Temporal position prior ──────────────────────────────────────────────
    total_duration = caption["duration"] or 1.0
    slide_idx = np.arange(N, dtype=np.float64)
    for t, seg in enumerate(segments):
        t_mid = (seg["start"] + seg["end"]) / 2.0
        expected_s = (t_mid / total_duration) * (N - 1)
        prior = -0.5 * ((slide_idx - expected_s) / PRIOR_SIGMA) ** 2
        log_ll[t] += prior

    # ── Off-slide detection ──────────────────────────────────────────────────
    raw_max_sim = log_ll_raw.max(axis=1)
    off_slide_mask = (raw_max_sim < OFF_SLIDE_THRESHOLD).tolist()
    n_off = sum(off_slide_mask)
    if n_off:
        print(f"  Off-slide segments: {n_off}/{T} (raw cosine < {OFF_SLIDE_THRESHOLD})")

    # ── Viterbi with enhanced temporal heuristics ────────────────────────────
    print("  Running Viterbi smoothing (enhanced temporal heuristics)...")
    slide_path_idx = viterbi_smooth_fast(log_ll)

    # ── Load slide text for labels ───────────────────────────────────────────
    slides = load_slides(slide_path)

    # ── Build output ─────────────────────────────────────────────────────────
    aligned_segments = []
    for i, seg in enumerate(segments):
        si = slide_path_idx[i]
        off = off_slide_mask[i]
        raw_sim = round(float(log_ll_raw[i, si]), 4)
        aligned_segments.append({
            "id": seg["id"],
            "start": seg["start"],
            "end": seg["end"],
            "text": seg["text"],
            "slide": None if off else si + 1,
            "slide_label": None if off else slides[si].label if si < len(slides) else f"Slide {si+1}",
            "similarity": raw_sim,
            "off_slide": off,
        })

    timeline = build_timeline(segments, slide_path_idx, slides, off_slide_mask)

    result = {
        "lecture": caption_path.stem,
        "slide_file": slide_path.name,
        "total_slides": N,
        "total_segments": T,
        "off_slide_count": n_off,
        "duration": caption["duration"],
        "language": caption.get("language", ""),
        "embed_model": "jina-embeddings-v4",
        "context_sec": CONTEXT_SEC,
        "segments": aligned_segments,
        "timeline": timeline,
    }

    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / f"{caption_path.stem}.json"
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    covered = len({s["slide"] for s in aligned_segments if s["slide"] is not None})
    print(f"  Covered {covered}/{N} slides across {len(timeline)} intervals")
    print(f"  Saved → {out_file}")

    # Cleanup temporary images
    import shutil
    shutil.rmtree(tmp_img_dir, ignore_errors=True)

    return out_file


# ── Auto-discovery ────────────────────────────────────────────────────────────

def _candidate_slides(course_dir: Path) -> list[Path]:
    """All PDF/PPTX/DOCX under [course]/materials/."""
    exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
    return [p for p in (course_dir / "materials").rglob("*") if p.suffix.lower() in exts]


def _lec_num(path: Path) -> int | None:
    """Extract lecture number from a filename, e.g. L02-foo → 2. Returns None if not found.

    The negative lookbehind prevents matching 'l' inside words like 'tutorial02'.
    """
    m = re.search(r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_]?0*(\d+)", path.stem)
    return int(m.group(1)) if m else None


def _name_similarity(a: str, b: str) -> float:
    """Rough token-overlap score between two filenames (lowercase, no ext)."""
    ta = set(a.lower().replace("-", " ").replace("_", " ").split())
    tb = set(b.lower().replace("-", " ").replace("_", " ").split())
    if not ta or not tb:
        return 0.0
    return len(ta & tb) / len(ta | tb)


def find_best_slide(caption_path: Path, slide_candidates: list[Path]) -> Path | None:
    """Return the slide file most likely corresponding to the caption (single-file fallback)."""
    stem = caption_path.stem
    best_score, best = 0.0, None
    for sp in slide_candidates:
        score = _name_similarity(stem, sp.stem)
        if score > best_score:
            best_score, best = score, sp
    return best if best_score > 0.05 else None


def _find_best_slide_group(
    caption_path: Path,
    slides_by_num: dict[int, list[Path]],
    all_slides: list[Path],
) -> list[Path]:
    """
    Return the best-matching group of slide files for a caption.

    Priority:
      1. If the caption filename contains a lecture number that exists in
         slides_by_num, return that group (exact number match).
      2. Otherwise, find the single best-matching slide by name similarity
         and return its group.
      3. If no group found, return the single best-matching slide as a list.
    """
    # Try to extract a lecture number directly from the caption name
    num = _lec_num(caption_path)
    if num is not None and num in slides_by_num:
        return sorted(slides_by_num[num])

    # Name-similarity fallback: find best matching slide then return its group
    best = find_best_slide(caption_path, all_slides)
    if best is None:
        return []
    best_num = _lec_num(best)
    if best_num is not None and best_num in slides_by_num:
        return sorted(slides_by_num[best_num])
    return [best]


def _sample_caption_text(caption_path: Path, max_segments: int = 40) -> str:
    """Return text sampled evenly across the whole transcript."""
    try:
        with open(caption_path, encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        return ""
    segs = data.get("segments", [])
    if not segs:
        return ""
    step = max(1, len(segs) // max_segments)
    return " ".join(
        s["text"].strip() for s in segs[::step][:max_segments] if s.get("text")
    )


def _sample_slide_text(slide_path: Path, max_words: int = 400) -> str:
    """Return text sampled evenly across all slides (no vision API)."""
    try:
        slides = load_slides(slide_path)   # describer=None → text only
    except Exception:
        return ""
    if not slides:
        return ""
    step = max(1, len(slides) // 15)
    sample = slides[::step][:15]
    words = " ".join(s.text for s in sample).split()
    return " ".join(words[:max_words])


def _content_match_slide_group(
    caption_path: Path,
    slides_by_num: dict[int, list[Path]],
    all_slides: list[Path],
    embedder,
    threshold: float = 0.20,
) -> list[Path]:
    """
    Content-based fallback for when name matching fails.

    Embeds a short sample of the caption transcript and a short sample of each
    candidate slide file, then returns the slide-file group whose text is most
    similar to the caption.  Returns [] when the best cosine similarity is below
    *threshold* (prevents spurious matches when no slide is truly related).
    """
    caption_text = _sample_caption_text(caption_path)
    if not caption_text.strip():
        return []

    # Build per-slide-file samples (skip files that yield no text)
    valid: list[Path] = []
    slide_texts: list[str] = []
    for sp in all_slides:
        t = _sample_slide_text(sp)
        if t.strip():
            valid.append(sp)
            slide_texts.append(t)

    if not valid:
        return []

    # Embed everything in one batch for efficiency
    all_texts = [caption_text] + slide_texts
    embs       = embed_texts(embedder, all_texts)   # already L2-normalised
    cap_emb    = embs[0:1]                          # (1, D)
    slide_embs = embs[1:]                           # (M, D)
    sims       = (cap_emb @ slide_embs.T)[0]        # (M,) cosine similarities

    best_idx = int(np.argmax(sims))
    best_sim  = float(sims[best_idx])

    if best_sim < threshold:
        print(f"  [content-match] {caption_path.name}: best sim {best_sim:.3f} "
              f"< {threshold} — no confident match found")
        return []

    best = valid[best_idx]
    print(f"  [content-match] {caption_path.name} → {best.name}  (sim={best_sim:.3f})")

    # Expand to the full group for this slide file
    best_num = _lec_num(best)
    if best_num is not None and best_num in slides_by_num:
        return sorted(slides_by_num[best_num])
    return [best]


# ── Embedding-based video↔slide matching ─────────────────────────────────────

MATCH_MODELS = {
    "bge-m3":      "BAAI/bge-m3",
    "jina":        "jina-embeddings-v4",   # via Jina API
    "google":      "google-embedding-v2",  # placeholder — not yet implemented
    "mpnet":       "all-mpnet-base-v2",    # local fallback
}


def suggest_matches(
    course_dir: Path,
    model: str = "bge-m3",
) -> dict[str, str]:
    """Match each caption to its best slide file using embedding similarity.

    Embeds a text sample from each transcript and each slide file, then
    finds the highest cosine-similarity pair for every caption.

    Args:
        course_dir: path containing captions/ and materials/
        model: embedding model key — "bge-m3" (default), "jina", "mpnet"

    Returns:
        {caption_stem: "materials/relative/path.pdf", ...}
    """
    captions_dir = course_dir / "captions"
    captions = sorted(captions_dir.glob("*.json")) if captions_dir.exists() else []
    all_slides = _candidate_slides(course_dir)

    if not captions or not all_slides:
        return {}

    print(f"  [match] Embedding-based matching with model={model}")
    print(f"  [match] {len(captions)} caption(s), {len(all_slides)} slide file(s)")

    # ── Sample text from each source ─────────────────────────────────────────
    cap_texts: list[str] = []
    cap_stems: list[str] = []
    for cap in captions:
        t = _sample_caption_text(cap)
        if t.strip():
            cap_texts.append(t)
            cap_stems.append(cap.stem)

    slide_texts: list[str] = []
    slide_paths: list[Path] = []
    for sp in all_slides:
        t = _sample_slide_text(sp)
        if t.strip():
            slide_texts.append(t)
            slide_paths.append(sp)

    if not cap_texts:
        print(f"  [match] No caption text extracted from {len(captions)} file(s) — check captions dir")
        return {}
    if not slide_texts:
        print(f"  [match] No slide text extracted from {len(all_slides)} file(s) — check materials dir")
        print(f"  [match] Slide paths tried: {[s.name for s in all_slides[:5]]}")
        return {}

    print(f"  [match] Extracted text: {len(cap_texts)} caption(s), {len(slide_texts)} slide(s)")

    # ── Embed with selected model ────────────────────────────────────────────
    if model == "jina":
        cap_embs = embed_texts_jina(cap_texts, desc="  [match] captions")
        slide_embs = embed_texts_jina(slide_texts, desc="  [match] slides")
        if cap_embs is None or slide_embs is None:
            print("  [match] Jina API failed — falling back to heuristics")
            return {}
    else:
        # Local model: bge-m3 or mpnet
        model_name = MATCH_MODELS.get(model, MATCH_MODELS["bge-m3"])
        try:
            from sentence_transformers import SentenceTransformer
            import torch
            device = "cuda" if torch.cuda.is_available() else "cpu"
            print(f"  [match] Loading {model_name} on {device}...")
            embedder = SentenceTransformer(model_name, device=device)
            all_texts = cap_texts + slide_texts
            embs = embedder.encode(
                all_texts,
                batch_size=32,
                normalize_embeddings=True,
                show_progress_bar=False,
                convert_to_numpy=True,
            ).astype(np.float32)
            cap_embs = embs[:len(cap_texts)]
            slide_embs = embs[len(cap_texts):]
            print(f"  [match] Embedded {len(cap_texts)} + {len(slide_texts)} texts")
        except Exception as e:
            import traceback
            print(f"  [match] Model load failed: {e}")
            traceback.print_exc()
            print("  [match] Falling back to heuristics")
            return {}

    # ── Cosine similarity → best match per caption ───────────────────────────
    sims = cap_embs @ slide_embs.T  # (n_caps, n_slides)

    # Build a lookup for "upgrade" candidates:
    # If we match "Lecture1 Intro.pdf", prefer "Lecture1 Intro With notes.pdf"
    _upgrade: dict[str, Path] = {}  # base_stem → best version path
    for sp in all_slides:
        base = re.sub(r"\s+(With notes|Review|ann\d+)\b", "", sp.stem,
                       flags=re.IGNORECASE).strip()
        existing = _upgrade.get(base)
        if existing is None:
            _upgrade[base] = sp
        else:
            # Prefer: "With notes" > "Review" > annotated > plain
            def _rank(p: Path) -> int:
                n = p.name.lower()
                if "with notes" in n: return 4
                if "review" in n:     return 3
                if re.search(r"ann\d", n): return 2
                return 1
            if _rank(sp) > _rank(existing):
                _upgrade[base] = sp

    result: dict[str, str] = {}
    for i, stem in enumerate(cap_stems):
        best_j = int(np.argmax(sims[i]))
        best_sim = float(sims[i, best_j])
        best_path = slide_paths[best_j]

        # Try to upgrade to a better version of the same lecture
        base = re.sub(r"\s+(With notes|Review|ann\d+)\b", "", best_path.stem,
                       flags=re.IGNORECASE).strip()
        upgraded = _upgrade.get(base, best_path)
        if upgraded != best_path:
            print(f"  [match] {stem} → {upgraded.name}  (sim={best_sim:.3f}, "
                  f"upgraded from {best_path.name})")
            best_path = upgraded
        else:
            print(f"  [match] {stem} → {best_path.name}  (sim={best_sim:.3f})")

        rel = str(best_path.relative_to(course_dir))
        if best_sim > 0.15:
            result[stem] = rel

    return result


def _load_mapping(mapping_path: Path, course_dir: Path) -> dict[str, list[Path]]:
    """Load a user-supplied video↔slide mapping JSON.

    Format: {"caption_stem": ["path/to/slide.pdf", ...], ...}
    Paths are relative to course_dir or absolute.
    Returns {caption_stem: [resolved_Path, ...]}.
    """
    with open(mapping_path, encoding="utf-8") as f:
        raw = json.load(f)
    mapping: dict[str, list[Path]] = {}
    for cap_stem, slide_list in raw.items():
        resolved = []
        for sp in slide_list:
            p = Path(sp)
            if not p.is_absolute():
                p = course_dir / p
            if p.exists():
                resolved.append(p)
            else:
                print(f"  [warn] Mapped slide not found: {sp}")
        if resolved:
            mapping[cap_stem] = resolved
    return mapping


def process_course(course_id: int | str, use_jina: bool = False,
                   mapping_path: Path | None = None) -> None:
    course_dir  = COURSE_DATA_DIR / str(course_id)
    captions_dir = course_dir / "captions"
    print(f"Course dir : {course_dir}", flush=True)
    print(f"Captions   : {captions_dir}", flush=True)

    captions = sorted(captions_dir.glob("*.json")) if captions_dir.exists() else []
    all_slides = _candidate_slides(course_dir)
    out_dir    = course_dir / "alignment"

    if not captions:
        print(f"[warn] No captions found in {captions_dir}")
        print("       Run 'Transcribe' first before aligning.")
        return
    if not all_slides:
        print(f"[warn] No slide files found under {course_dir / 'materials'}")
        print("       Run 'Download materials' first.")
        return

    print(f"Found {len(captions)} caption(s), {len(all_slides)} slide file(s).", flush=True)

    # ── Load user-supplied mapping if provided ────────────────────────────────
    user_mapping: dict[str, list[Path]] = {}
    if mapping_path and mapping_path.exists():
        user_mapping = _load_mapping(mapping_path, course_dir)
        print(f"  User mapping: {len(user_mapping)} video→slide pair(s) loaded")
        for cap_stem, slides in user_mapping.items():
            print(f"    {cap_stem} → {[s.name for s in slides]}")

    # Check Jina API availability if requested
    if use_jina:
        jina_key = _get_jina_key()
        if jina_key:
            print("  [jina] Jina API key found — using multimodal alignment")
        else:
            print("  [jina] No Jina API key — falling back to text-based alignment")
            use_jina = False

    # Group slide files by lecture number
    slides_by_num: dict[int, list[Path]] = defaultdict(list)
    for sp in all_slides:
        num = _lec_num(sp)
        if num is not None:
            slides_by_num[num].append(sp)

    embedder = None if use_jina else None  # lazy-load only when needed

    for cap in captions:
        # Skip captions flagged as low-quality (wrong/empty recordings)
        try:
            with open(cap, encoding="utf-8") as _f:
                _meta = json.load(_f)
            if _meta.get("quality") == "low":
                print(f"  [skip] Low-quality transcript (wrong/empty recording): {cap.name}")
                continue
        except Exception:
            pass   # can't read → proceed and let align() handle it

        # ── Priority 1: user-supplied mapping ────────────────────────────────
        slide_group: list[Path] = []
        if cap.stem in user_mapping:
            slide_group = user_mapping[cap.stem]
            print(f"  [mapping] {cap.name} → {[s.name for s in slide_group]}")
        else:
            # ── Priority 2: automatic name/number matching ───────────────────
            slide_group = _find_best_slide_group(cap, slides_by_num, all_slides)
            if not slide_group:
                # ── Priority 3: content embedding fallback ───────────────────
                if embedder is None:
                    embedder = get_embedder()
                slide_group = _content_match_slide_group(
                    cap, slides_by_num, all_slides, embedder
                )
            if not slide_group:
                print(f"  [warn] No matching slides for {cap.name} — skipping")
                continue

        # Check if all output files for this group already exist
        if len(slide_group) == 1:
            out_file = out_dir / f"{cap.stem}.json"   # legacy single-file naming
            if out_file.exists():
                print(f"  [skip] Already aligned: {cap.stem}")
                continue
        else:
            if all((out_dir / f"{sp.stem}.json").exists() for sp in slide_group):
                print(f"  [skip] Already aligned: {[sp.name for sp in slide_group]}")
                continue

        # Use Jina multimodal for single PDF slides if available
        if use_jina and len(slide_group) == 1 and slide_group[0].suffix.lower() == ".pdf":
            result = align_multimodal(cap, slide_group[0], out_dir)
            if result:
                continue
            # Fall back to text-based on failure
            print("  [jina] Falling back to text-based alignment")

        if embedder is None:
            embedder = get_embedder()
        align_multi_slides(cap, slide_group, out_dir, embedder=embedder)


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Semantic caption↔slide alignment")
    parser.add_argument("--caption", metavar="PATH",
                        help="Whisper caption JSON")
    parser.add_argument("--slides",  metavar="PATH", nargs="+",
                        help="One or more slide files (PDF/PPTX/DOCX). "
                             "Multiple files are concatenated in order.")
    parser.add_argument("--course",  metavar="ID",
                        help="Auto-process a course folder")
    parser.add_argument("--out",     metavar="DIR",
                        help="Output directory (default: [course]/alignment)")
    parser.add_argument("--jina",    action="store_true",
                        help="Use Jina Embeddings v4 for multimodal (image+text) alignment")
    parser.add_argument("--mapping", metavar="JSON",
                        help="User-supplied video↔slide mapping JSON file. "
                             "Format: {\"caption_stem\": [\"slide_path\", ...]}. "
                             "Unmapped captions fall back to auto-discovery.")
    parser.add_argument("--suggest-matches", action="store_true",
                        help="Suggest video↔slide matches using embedding similarity "
                             "(outputs JSON to stdout). Does NOT run alignment.")
    parser.add_argument("--match-model", metavar="MODEL", default="bge-m3",
                        choices=["bge-m3", "jina", "mpnet"],
                        help="Embedding model for --suggest-matches "
                             "(default: bge-m3, alternatives: jina, mpnet)")
    args = parser.parse_args()

    if args.suggest_matches:
        if not args.course:
            parser.error("--suggest-matches requires --course ID")
        course_dir = COURSE_DATA_DIR / str(args.course)
        matches = suggest_matches(course_dir, model=args.match_model)
        # Output as JSON for IPC consumption
        print("\n__MATCH_RESULT__")
        print(json.dumps(matches, indent=2))
        return

    if args.course:
        mapping = Path(args.mapping) if args.mapping else None
        process_course(args.course, use_jina=args.jina, mapping_path=mapping)
        return

    if not args.caption or not args.slides:
        parser.error("Provide both --caption and --slides, or use --course ID")

    cap_path    = Path(args.caption)
    slide_paths = [Path(p) for p in args.slides]

    if not cap_path.exists():
        print(f"[error] Caption not found: {cap_path}"); sys.exit(1)
    for sp in slide_paths:
        if not sp.exists():
            print(f"[error] Slides not found: {sp}"); sys.exit(1)

    # Infer output dir from caption path: two levels up → alignment/
    out_dir = Path(args.out) if args.out else cap_path.parent.parent / "alignment"

    embedder = get_embedder()
    align_multi_slides(cap_path, slide_paths, out_dir, embedder=embedder)


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n[info] Interrupted by user.")
        sys.exit(0)
    except SystemExit:
        raise
    except Exception as _exc:
        import traceback
        print(f"\n[error] Unexpected error: {_exc}")
        traceback.print_exc()
        sys.exit(1)
