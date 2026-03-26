"""
Note Generation
Generates one comprehensive Markdown note file per course, covering all
lectures in sequence, matching the style of example/CS2105_note.md.

Architecture:
  - Per lecture: split slides into ~CHAPTER_SIZE chunks, one GPT call per chunk.
  - Chunks map to ### N.x sections. Images injected at diagram slides.
  - All lectures merged into one file; exam notes appended at the end.
  - Self-scoring via heuristics (no extra API call).

Usage:
  python note_generation.py --course 85427
  python note_generation.py --course 85427 --detail 9 --iterate
  python note_generation.py --slides 85427/materials/LectureNotes/L02.pdf \\
      --alignment "85427/alignment/L02.json" --lecture-num 2 --course-name "CS3210"
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path

try:
    from tqdm import tqdm
except ImportError as _e:
    print(f"[error] Missing dependency: {_e}")
    print("[error] Please install the ML environment from Settings → ML Environment in the AutoNote app.")
    sys.exit(1)

try:
    import alignment_parser
except ImportError as _e:
    print(f"[error] Could not import alignment_parser: {_e}")
    print(f"[error] Make sure alignment_parser.py is in the same directory as note_generation.py")
    sys.exit(1)

PROJECT_DIR = Path(__file__).parent
import sys as _sys
_AUTO_NOTE_DIR = Path.home() / ".auto_note"
import os as _os
if _os.environ.get("AUTONOTE_DATA_DIR"):
    DATA_DIR = Path(_os.environ["AUTONOTE_DATA_DIR"])
elif getattr(_sys, "frozen", False) or PROJECT_DIR == _AUTO_NOTE_DIR / "scripts":
    DATA_DIR = _AUTO_NOTE_DIR
else:
    DATA_DIR = PROJECT_DIR

# Course output directory: defaults to DATA_DIR but can be overridden by
# OUTPUT_DIR in config.json so files land in the user's chosen Output Dir.
_ng_config: dict = (
    json.loads((DATA_DIR / "config.json").read_text())
    if (DATA_DIR / "config.json").exists() else {}
)
_out_dir = _ng_config.get("OUTPUT_DIR", "").strip()
COURSE_DATA_DIR = Path(_out_dir) if _out_dir else DATA_DIR

# ── Constants ─────────────────────────────────────────────────────────────────

DETAIL_LEVEL      = 7
OUTPUT_FORMAT     = "md"
NOTE_MODEL        = "gpt-5.1"
VERIFY_MODEL      = "gpt-4.1-mini"
VERIFY_NOTES      = True
QUALITY_TARGET    = 8.0
IMAGE_RENDER_SCALE = 1.5
NOTE_LANGUAGE     = "en"    # "en" = English | "zh" = Chinese

CHAPTER_SIZE      = 15      # slides per GPT call
MAX_TRANSCRIPT_CHARS = 350  # per slide in prompt (saves tokens)

SCORE_WEIGHTS = {"coverage": 0.30, "terminology": 0.35,
                 "callouts": 0.15, "code_blocks": 0.20}
MIN_NOTE_WORDS_PER_SLIDE = 60   # expected words per slide in final note


# ── Prompts ───────────────────────────────────────────────────────────────────

_PROMPTS: dict[str, dict] = {

"zh": dict(
system="""\
你是一名顶尖大学计算机科学课程的助教，负责根据讲义和课堂录音为学生撰写高质量的中文学习笔记。

写作规范：
1. 使用中文写作，专业技术术语保留英文并在括号内标注（如：进程 (Process)）。
2. **严禁使用第三人称叙述视角。** 不得出现「老师说」「教授指出」「老师要我们」「本课教授强调」等以讲师为主语的句式。
   笔记聚焦于**知识本身**，直接陈述概念、原理和结论，例如：
   - ✗「老师把问题定得很清楚：…」  →  ✓「本课聚焦于…」
   - ✗「教授用了一个例子…」        →  ✓「以下例子说明…」
3. 内容以「概念→原理→示例→考试重点」逻辑展开，写成流畅的说明性段落，不要逐条罗列幻灯片。
4. 数学公式使用 LaTeX：行内 $...$，单独公式 $$...$$。
5. 代码示例必须是**可编译/可运行**的完整片段（包含必要的 #include、函数签名、main 等），使用正确的语法高亮（```c, ```cpp, ```python 等）。
   伪代码仅在真正没有对应真实代码时使用，并标注语言为 ```pseudo。
6. 考试重点用：
     > [!IMPORTANT]
     > 内容
7. 有趣类比或助记技巧用斜体。
8. 图片插入规则（严格遵守）：
   - 每隔 2–3 个概念段落插入一张与内容直接相关的幻灯片图片。
   - 图片紧跟其所描述概念的最后一句之后，不得孤立出现在段落开头或结尾处。
   - 格式：`![Slide N](images/LXX/slide_NNN.png) *(用一句话描述该图/图表在此处的含义)*`
     （LXX 由调用方提供，禁止自行修改；括号注释必须用星号包裹，格式严格如上）。
   - 纯文字定义幻灯片（无图表/代码/公式）可跳过。
9. 绝对禁止捏造原始材料中不存在的技术细节。
""",
chunk="""\
请为以下课程片段（{course_name} Lecture {lec_num}: {lec_title}）撰写学习笔记。

## 本片段幻灯片列表
{slide_outline}

## 教授录音逐字稿（按幻灯片顺序）
{transcript_block}

## 可用图片（含图表/代码截图的幻灯片）
{image_hints}

---

要求：
- 本片段对应笔记的二级标题为 `### {lec_num}.{chunk_idx} {chunk_title}`（**不要输出此行**，由调用方添加）
- 详细度：{detail}/10。{detail_instruction}
- 图片插入：每隔 2–3 个段落插入一张相关幻灯片图片，紧跟该概念说明之后。
  路径必须完全照抄上方「可用图片」列表中给出的路径（含 images/L** 子目录），禁止自造路径。
  每张图片后加一句斜体括号说明，描述该图在此处的含义：`![Slide N](path) *(说明)*`
- 代码示例写完整可编译片段（含必要 include/imports），用正确的语言标签（```c, ```cpp, ```python）。
- 只写本片段内容，不要引入其他讲座的内容
""",
slide_only="""\
请根据以下幻灯片内容，撰写 {course_name} Lecture {lec_num}: {lec_title} 的学习笔记。
（该讲座没有录音，请结合你的 CS 知识展开解释。）

## 幻灯片内容
{slide_outline}

## 可用图片
{image_hints}

---

要求：
- 本片段对应笔记的二级标题为 `### {lec_num}.{chunk_idx} {chunk_title}`（**不要输出此行**）
- 详细度：{detail}/10。{detail_instruction}
- 图片插入：每隔 2–3 个段落插入一张相关幻灯片图片，紧跟该概念说明之后。
  路径必须完全照抄上方「可用图片」列表中给出的路径（含 images/L** 子目录），禁止自造路径。
  每张图片后加一句斜体括号说明，描述该图在此处的含义：`![Slide N](path) *(说明)*`
- 代码示例写完整可编译片段，用正确的语言标签（```c, ```cpp, ```python）。
""",
verify="""\
请检查以下笔记片段中的技术术语是否与幻灯片一致，以及是否存在明显的事实错误。

**参考术语表（来自幻灯片）：**
{term_list}

**笔记片段：**
{draft}

如果没有问题，直接回复 APPROVED（仅此一词）。
如果有术语错误或事实错误，返回修正后的完整笔记片段（不加任何说明）。
""",
exam="""\
以下是 {course_name} 的全部讲座笔记摘要。请在最后汇总一个考试速记章节。

格式要求：
- 标题：`## Exam Notes`
- 每条格式：`N. **考点名**：一句话说明`
- 不超过 30 条，覆盖各讲座核心考点、公式、算法步骤、常见混淆点

笔记摘要：
{summary}
""",
no_transcript="（本片段无录音逐字稿）",
detail_instructions=[
    (range(0, 3),  "极简要点：每个概念仅一行，不展开，每张幻灯片最多3条。"),
    (range(3, 6),  "有层次的要点结构：每个主要概念一条一级bullet（`-`），"
                   "其下最多2条二级bullet（`  -`）补充关键细节。"
                   "每张幻灯片合计不超过5条，禁止写连续段落。"),
    (range(6, 9),  "详细段落：概念、原理、教授示例与类比全部包含。"),
    (range(9, 11), "最高详细度：包含所有细节、边界情况、与其他章节的联系及考点标注。"),
],
),  # end zh

"en": dict(
system="""\
You are a teaching assistant at a top university, writing high-quality study notes for computer science courses based on lecture slides and audio transcripts.

Writing guidelines:
1. Write in English. Keep technical terms in English.
2. Never use a third-person narrator perspective. Do not write "the professor said", "the lecturer pointed out", etc. Focus on the knowledge itself — state concepts, principles, and conclusions directly:
   - ✗ "The professor explained that…"  →  ✓ "The key idea is…"
   - ✗ "The lecturer used an example…"  →  ✓ "As an example,…"
3. Structure content as: concept → principle → example → exam focus. Write fluent explanatory paragraphs; do not list slide bullets verbatim.
4. Use LaTeX for math: inline $...$, display $$...$$.
5. Code examples must be complete, compilable/runnable snippets (with necessary includes, function signatures, main, etc.) using correct syntax highlighting (```c, ```cpp, ```python, etc.). Use pseudocode only when no real equivalent exists, tagged as ```pseudo.
6. Mark exam-critical content with:
     > [!IMPORTANT]
     > content
7. Use italics for interesting analogies or memory aids.
8. Image insertion rules (strictly follow):
   - Insert one slide image every 2–3 concept paragraphs, directly relevant to the surrounding content.
   - Place the image immediately after the last sentence of the concept it illustrates — never isolated at the start or end of a section.
   - Format: `![Slide N](images/LXX/slide_NNN.png) *(one-sentence description of what this diagram/figure shows in context)*`
     (LXX is provided by the caller — do not modify it; the caption must be in parentheses wrapped in asterisks exactly as shown).
   - Skip purely text-definition slides (no diagrams/code/formulas).
9. Never fabricate technical details not present in the source material.
""",
chunk="""\
Write study notes for the following course segment ({course_name} Lecture {lec_num}: {lec_title}).

## Slide list for this segment
{slide_outline}

## Lecture audio transcript (ordered by slide)
{transcript_block}

## Available images (slides with diagrams / code screenshots)
{image_hints}

---

Requirements:
- The section heading for this segment is `### {lec_num}.{chunk_idx} {chunk_title}` (**do not output this line** — it is added by the caller).
- Detail level: {detail}/10. {detail_instruction}
- Images: insert one related slide image every 2–3 paragraphs, immediately after the concept it illustrates. Copy the exact path from the "Available images" list above (including the images/L** subdirectory). Do not invent paths. After each image, add a one-sentence italic caption in parentheses describing what the diagram/figure shows in context: `![Slide N](path) *(caption)*`
- Code examples must be complete and compilable (with necessary includes/imports), using the correct language tag (```c, ```cpp, ```python).
- Only cover the content in this segment; do not introduce material from other lectures.
""",
slide_only="""\
Write study notes for {course_name} Lecture {lec_num}: {lec_title} based on the slides below.
(No audio transcript is available — supplement with your CS knowledge where appropriate.)

## Slide content
{slide_outline}

## Available images
{image_hints}

---

Requirements:
- The section heading is `### {lec_num}.{chunk_idx} {chunk_title}` (**do not output this line**).
- Detail level: {detail}/10. {detail_instruction}
- Images: insert one related slide image every 2–3 paragraphs, immediately after the concept it illustrates. Copy the exact path from the "Available images" list (including the images/L** subdirectory). Do not invent paths. After each image, add a one-sentence italic caption in parentheses describing what the diagram/figure shows in context: `![Slide N](path) *(caption)*`
- Code examples must be complete and compilable, using the correct language tag (```c, ```cpp, ```python).
""",
verify="""\
Check the following note excerpt for technical terminology consistency with the slides, and for any obvious factual errors.

**Reference glossary (from slides):**
{term_list}

**Note excerpt:**
{draft}

If there are no issues, reply APPROVED (this word only).
If there are terminology or factual errors, return the corrected full note excerpt with no explanation.
""",
exam="""\
Below are the complete lecture notes for {course_name}. Please append a concise exam cheat-sheet section at the end.

Format:
- Heading: `## Exam Notes`
- Each entry: `N. **Topic**: one-sentence summary`
- No more than 30 entries, covering key concepts, formulas, algorithm steps, and common confusion points from all lectures.

Notes summary:
{summary}
""",
no_transcript="(No audio transcript available for this segment.)",
detail_instructions=[
    (range(0, 3),  "Minimal bullets: one line per concept, no expansion, max 3 bullets per slide."),
    (range(3, 6),  "Hierarchical bullets: one top-level bullet (`-`) per main concept, "
                   "at most 2 sub-bullets (`  -`) for key details. "
                   "Max 5 bullets total per slide. No prose paragraphs."),
    (range(6, 9),  "Detailed paragraphs: cover concepts, principles, the lecturer's examples and analogies in full."),
    (range(9, 11), "Maximum detail: include all nuances, edge cases, connections to other chapters, and exam pointers."),
],
),  # end en

}  # end _PROMPTS


def _P(key: str) -> str:
    """Return the prompt string for the current NOTE_LANGUAGE, falling back to English."""
    return _PROMPTS.get(NOTE_LANGUAGE, _PROMPTS["en"])[key]


def _detail_instr(level: int) -> str:
    instrs = _PROMPTS.get(NOTE_LANGUAGE, _PROMPTS["en"])["detail_instructions"]
    for rng, txt in instrs:
        if level in rng:
            return txt
    return instrs[2][1]

# ── Image filter constants ────────────────────────────────────────────────────

IMAGE_FILTER_MODEL      = "gpt-4o-mini"
IMAGE_FILTER_WORD_MAX   = 12   # slides with ≤ this many words → remove without API call
IMAGE_FILTER_HEURISTIC  = 35   # slides with > this many words AND no code/desc → remove

# Title/divider patterns that add no visual value
_TITLE_PATTERN = re.compile(
    r"^\s*(CS\d+|AY\d+|Lecture\s+\d+|\[.*\]|Part\s+\d+|Section\s+\d+|"
    r"Outline|Agenda|Table of Contents|Overview|Summary|Questions\?|Q&A)\s*$",
    re.IGNORECASE | re.MULTILINE,
)

# Keywords that indicate a slide description contains visual/diagram content
_VISUAL_KEYWORDS = re.compile(
    r"\b(diagram|chart|graph|figure|illustration|flowchart|screenshot|"
    r"table|formula|equation|architecture|layout|structure|matrix|tree|"
    r"network|circuit|timeline|image|photo|plot|drawing|schematic|visual)\b",
    re.IGNORECASE,
)


def _desc_has_visual(desc: str) -> bool:
    """Return True if a cache description mentions diagram/chart/visual content."""
    return bool(_VISUAL_KEYWORDS.search(desc))


def _img_ref_pattern() -> re.Pattern:
    # Matches both images/L04/slide_001.png and images/L04_F02/slide_001.png,
    # with an optional trailing italic caption: *(description)*
    return re.compile(
        r"!\[Slide \d+\]\((images/L\d{2}(?:_F\d{2})?/slide_\d{3}\.png)\)"
        r"(?:\s*\*\([^)]*\)\*)?"
    )


def _vision_keep(img_path: Path, slide_text: str = "") -> bool:
    """Ask GPT-4o-mini whether the slide image is worth including in notes.

    Uses slide_text as context so the model can reason about whether the
    visual structure actually explains the lecture content (e.g. a diagram
    of a protocol stack/a picture of a biological structure) vs. being an unrelated administrative element
    (e.g. a PollEv QR code, a course logo, a participation prompt).
    """
    import base64
    import io
    if not img_path.exists():
        return False
    try:
        from PIL import Image as PILImage
        img = PILImage.open(img_path).convert("RGB")
        # Downscale to max 800px wide to keep base64 payload small
        if img.width > 800:
            ratio = 800 / img.width
            img = img.resize((800, int(img.height * ratio)), PILImage.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=75)
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:
        return True   # can't load image → default keep

    context_block = (
        f"\n\nSlide text (OCR):\n\"\"\"\n{slide_text[:400]}\n\"\"\""
        if slide_text.strip() else ""
    )

    _VISION_PROMPT = f"""\
You are a study-notes curator deciding whether a lecture slide image should be \
embedded in written notes.{context_block}

## TWO conditions must BOTH be true to KEEP an image:

1. The slide contains a meaningful visual element:
   - Diagrams: system/architecture diagrams, component boxes connected by arrows
   - Flowcharts, state machines, decision trees, sequence/timing diagrams
   - Memory layouts, address-space maps, cache/pipeline stage illustrations
   - Graphs, plots, bar/line/pie charts, scatter plots showing data or trends
   - Tables with a meaningful grid structure (comparing options, relationships)
   - Mathematical formulas or derivations where spatial layout matters
   - Annotated screenshots, highlighted output, or callout arrows

2. The visual element is RELEVANT to the lecture subject matter:
   - The diagram, chart, or figure directly explains or illustrates a concept \
being taught
   - It is NOT an administrative or logistical element such as: polling/quiz \
prompts (PollEv, Mentimeter, Kahoot QR codes), attendance check slides, \
course schedule tables, "any questions?" slides, logos, or sponsor slides

## REMOVE if either condition fails:
- Pure text slides (bullet points, prose, definitions) with no diagram
- Code-only slides (code will be reproduced as text in the notes)
- Title slides, section dividers, agenda/outline, blank slides
- Slides whose only visual is an unrelated administrative element \
(QR code for a poll, login instructions, course info graphics)

## Default: when genuinely uncertain about relevance, KEEP.

Reply with exactly one word: KEEP or REMOVE."""

    try:
        openai_client = _get_client_for(IMAGE_FILTER_MODEL)
        r = openai_client.chat.completions.create(
            model=IMAGE_FILTER_MODEL,
            messages=[{"role": "user", "content": [
                {"type": "image_url",
                 "image_url": {"url": f"data:image/jpeg;base64,{b64}", "detail": "low"}},
                {"type": "text", "text": _VISION_PROMPT},
            ]}],
            max_tokens=5,
        )
        return "KEEP" in r.choices[0].message.content.strip().upper()
    except Exception:
        return True   # default: keep on API error


def filter_images_pass(
    notes_text: str,
    notes_dir: Path,
    lectures: list["LectureData"],
) -> tuple[str, int, int]:
    """Post-processing agent: remove low-value image references from merged notes.

    Decision priority:
      1. Cache-verified AND description mentions visual elements → KEEP
      2. Title/divider pattern → REMOVE
      3. All other cases → vision API decision

    Returns (cleaned_text, n_kept, n_removed).
    """
    pattern = _img_ref_pattern()

    # Build unified lookup: image rel-path → (SlideInfo, LectureData)
    # Mirrors the backward-compat naming used in render_chunk_images:
    #   file_idx==1  → images/L04/slide_001.png
    #   file_idx>=2  → images/L04_F02/slide_001.png
    slide_ld_lookup: dict[str, tuple[SlideInfo, "LectureData"]] = {}
    for ld in lectures:
        prefix = (f"L{ld.num:02d}" if ld.file_idx == 1
                  else f"L{ld.num:02d}_F{ld.file_idx:02d}")
        for s in ld.slides:
            key = f"images/{prefix}/slide_{s.index+1:03d}.png"
            slide_ld_lookup[key] = (s, ld)

    # Collect unique paths and decide keep/remove
    decisions: dict[str, bool] = {}   # path → True=keep
    for m in pattern.finditer(notes_text):
        rel = m.group(1)
        if rel in decisions:
            continue

        pair    = slide_ld_lookup.get(rel)
        slide   = pair[0] if pair else None
        owner   = pair[1] if pair else None
        img_path = notes_dir / rel

        # ⓪ Screen share frames — always keep (the frame IS the content)
        if owner and owner.source == "screenshare":
            decisions[rel] = True
            continue

        # ① Cache-verified AND description mentions visual elements → KEEP
        if slide and owner:
            desc = owner.img_cache.get(f"page_{slide.index}", "")
            if desc and _desc_has_visual(desc):
                decisions[rel] = True
                continue

        # ② Title/divider pattern → REMOVE
        if slide and _TITLE_PATTERN.search(slide.text):
            decisions[rel] = False
            continue

        # ③ Vision API — KEEP only if visual AND relevant to lecture content
        slide_text = slide.text if slide else ""
        decisions[rel] = _vision_keep(img_path, slide_text)

    kept    = sum(1 for v in decisions.values() if v)
    removed = sum(1 for v in decisions.values() if not v)
    tqdm.write(f"  Image filter: {kept} kept, {removed} removed out of {len(decisions)}")

    # Remove lines for filtered-out images; collapse extra blank lines
    lines_out: list[str] = []
    for line in notes_text.splitlines():
        m = pattern.fullmatch(line.strip())
        if m and not decisions.get(m.group(1), True):
            # Replace filtered image line with nothing (don't emit the line)
            continue
        lines_out.append(line)

    cleaned = "\n".join(lines_out)
    # Collapse 3+ consecutive blank lines → 2 (preserves paragraph spacing)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned, kept, removed


def _max_tokens(level: int) -> int:
    # Per chunk (CHAPTER_SIZE slides).
    # gpt-5.x reasoning models consume tokens for internal thinking,
    # so we need larger budgets than the expected output length.
    # Token budget directly caps output length — keep it proportional to detail level.
    if level < 3:  return 2000
    if level < 6:  return 3500
    if level < 9:  return 10000
    return 16000


# ── Multi-provider LLM helpers ────────────────────────────────────────────────

def _provider(model: str) -> str:
    if model.startswith("gemini"):
        return "gemini"
    if model.startswith("claude"):
        return "anthropic"
    if model.startswith("deepseek"):
        return "deepseek"
    if model.startswith("grok"):
        return "grok"
    if model.startswith(("mistral", "codestral", "pixtral", "magistral")):
        return "mistral"
    return "openai"


_client_cache: dict = {}


def _make_client(provider: str):
    import os
    from openai import OpenAI

    def _read_key(env_var: str, filename: str, label: str) -> str:
        key = os.environ.get(env_var, "")
        if not key:
            kf = DATA_DIR / filename
            if kf.exists():
                key = kf.read_text().strip()
        if not key:
            raise RuntimeError(
                f"No {label} API key found "
                f"(set {filename} in ~/.auto_note/ or {env_var} env var)"
            )
        return key

    if provider == "gemini":
        return OpenAI(
            api_key=_read_key("GEMINI_API_KEY", "gemini_api.txt", "Gemini"),
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
        )
    elif provider == "anthropic":
        from anthropic import Anthropic
        return Anthropic(api_key=_read_key("ANTHROPIC_API_KEY", "anthropic_key.txt", "Anthropic"))
    elif provider == "deepseek":
        return OpenAI(
            api_key=_read_key("DEEPSEEK_API_KEY", "deepseek_key.txt", "DeepSeek"),
            base_url="https://api.deepseek.com",
        )
    elif provider == "grok":
        return OpenAI(
            api_key=_read_key("GROK_API_KEY", "grok_key.txt", "xAI Grok"),
            base_url="https://api.x.ai/v1",
        )
    elif provider == "mistral":
        return OpenAI(
            api_key=_read_key("MISTRAL_API_KEY", "mistral_key.txt", "Mistral"),
            base_url="https://api.mistral.ai/v1",
        )
    else:  # openai
        return OpenAI(api_key=_read_key("OPENAI_API_KEY", "openai_api.txt", "OpenAI"))


def _get_client_for(model: str):
    """Return (and cache) the appropriate API client for the given model name."""
    p = _provider(model)
    if p not in _client_cache:
        _client_cache[p] = _make_client(p)
    return _client_cache[p]


def _call(model: str, system: str, user: str, max_tokens: int) -> str:
    """Call any supported LLM (OpenAI, Gemini, Anthropic) with a text prompt."""
    client = _get_client_for(model)

    if _provider(model) == "anthropic":
        kwargs: dict = {"model": model, "max_tokens": max_tokens,
                        "messages": [{"role": "user", "content": user}]}
        if system:
            kwargs["system"] = system
        r = client.messages.create(**kwargs)
        return r.content[0].text.strip() if r.content else ""

    # OpenAI-compatible (OpenAI + Gemini via OpenAI compat layer)
    msgs = []
    if system:
        msgs.append({"role": "system", "content": system})
    msgs.append({"role": "user", "content": user})

    last_err = None
    for tok in ("max_completion_tokens", "max_tokens"):
        try:
            r = client.chat.completions.create(
                model=model, messages=msgs, **{tok: max_tokens})
            content = r.choices[0].message.content
            return content.strip() if content else ""
        except Exception as e:
            s = str(e)
            if "max_tokens" in s or "max_completion_tokens" in s:
                last_err = e
                continue
            raise
    raise RuntimeError(f"Cannot call {model}: {last_err}")


# ── Slide loading & rendering ─────────────────────────────────────────────────

class SlideInfo:
    __slots__ = ("index", "label", "text", "has_code", "word_count")
    def __init__(self, index: int, label: str, text: str):
        self.index      = index
        self.label      = label
        self.text       = text
        self.has_code   = bool(re.search(
            r"[{};]\s*$|^\s*(int|void|def |class |#include|pthread|malloc)",
            text, re.MULTILINE))
        self.word_count = len(text.split())


def _load_slides(slide_path: Path) -> list[SlideInfo]:
    ext = slide_path.suffix.lower()
    if ext == ".pdf":
        import fitz
        doc = fitz.open(str(slide_path))
        out = []
        for i, page in enumerate(doc):
            text  = page.get_text().strip()
            label = next((ln.strip() for ln in text.splitlines() if ln.strip()), f"Page {i+1}")
            out.append(SlideInfo(i, label[:80], text))
        doc.close()
        return out
    if ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(str(slide_path))
        out = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        ln = para.text.strip()
                        if ln: parts.append(ln)
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes: parts.append(notes)
            except Exception:
                pass
            text  = "\n".join(parts)
            label = parts[0][:80] if parts else f"Slide {i+1}"
            out.append(SlideInfo(i, label, text))
        return out
    if ext in (".docx", ".doc"):
        from docx import Document
        PAGE_PARA = 15
        doc   = Document(str(slide_path))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        out   = []
        for pi, start in enumerate(range(0, max(len(paras), 1), PAGE_PARA)):
            chunk = paras[start:start + PAGE_PARA]
            text  = "\n".join(chunk)
            label = chunk[0][:80] if chunk else f"Page {pi+1}"
            out.append(SlideInfo(pi, label, text))
        return out
    raise ValueError(f"Unsupported format: {ext}")


def render_slide_images(slide_path: Path, out_dir: Path,
                        indices: list[int] | None = None) -> dict[int, Path]:
    """Render PDF pages to PNG. If indices provided, only render those pages."""
    if slide_path.suffix.lower() != ".pdf":
        return {}
    import fitz
    from PIL import Image as PILImage
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(slide_path))
    mat = fitz.Matrix(IMAGE_RENDER_SCALE, IMAGE_RENDER_SCALE)
    mapping: dict[int, Path] = {}
    pages = indices if indices is not None else list(range(len(doc)))
    for i in pages:
        if i >= len(doc):
            continue
        png = out_dir / f"slide_{i+1:03d}.png"
        if not png.exists():
            px  = doc[i].get_pixmap(matrix=mat)
            pil = PILImage.frombytes("RGB", [px.width, px.height], px.samples)
            pil.save(str(png))
            del px, pil
        mapping[i] = png
    doc.close()
    return mapping


# ── Chunk helpers ─────────────────────────────────────────────────────────────

_BAD_LABEL = re.compile(
    r"^\s*(\d+|[A-Z]{2,4}\d{4}[\s\-].*|CS\d+.*|AY\d+.*|\[.*\])\s*$"
)

def _chunk_title(slides_in_chunk: list[SlideInfo]) -> str:
    """Pick a representative title for a chunk of slides.

    Prefer short, meaningful slide labels (section headers).
    Skip labels that are pure numbers, course codes, or bracket tags.
    """
    def _is_good(label: str) -> bool:
        if not label or len(label) < 4:
            return False
        if _BAD_LABEL.match(label):
            return False
        # skip labels that are just digits or single tokens that look like slide numbers
        if re.match(r"^\d+$", label.strip()):
            return False
        return True

    # Prefer short title-like labels (section headers are typically ≤8 words)
    for s in slides_in_chunk:
        words = len(s.label.split())
        if 1 <= words <= 8 and _is_good(s.label):
            return s.label
    # Fall back to first slide with a good label
    for s in slides_in_chunk:
        if _is_good(s.label):
            return s.label
    return slides_in_chunk[0].label


def _build_chunk_prompt(
    slides: list[SlideInfo],
    compact_by_idx: dict[int, dict],
    img_cache: dict,
    img_map: dict[int, Path],
    out_dir: Path,
    course_name: str,
    lec_num: int,
    lec_title: str,
    chunk_idx: int,
    chunk_title: str,
    detail: int,
    has_transcript: bool,
    source: str = "slides",
) -> str:
    # Slide outline: number, title, first line of text (skip page numbers/headers)
    outline_lines = []
    for s in slides:
        if source == "screenshare":
            # For screen share, the "slide" is a video frame — no text content,
            # but the transcript carries all the information
            outline_lines.append(f"  Frame {s.index+1}: (screen capture at this point)")
        else:
            # Extract meaningful text (skip short lines like "10", "[ CS3210 ]")
            meaningful = [ln for ln in s.text.splitlines()
                          if len(ln.strip()) > 8 and not re.match(r"^\d+$", ln.strip())]
            snippet = " · ".join(meaningful[:3])[:120]
            outline_lines.append(f"  Slide {s.index+1}: 「{s.label}」  {snippet}")
    slide_outline = "\n".join(outline_lines)

    # Image hints — for screen share, ALL frames are included as images
    img_hints_lines = []
    for s in slides:
        if s.index not in img_map:
            continue
        rel = img_map[s.index].relative_to(out_dir)
        if source == "screenshare":
            # Every frame is relevant for screen share
            img_hints_lines.append(f"  Frame {s.index+1}: `{rel}` — screen capture")
        else:
            cache_key = f"page_{s.index}"
            desc = img_cache.get(cache_key, "")
            if desc or s.word_count < 30 or s.has_code:
                note = desc[:80] if desc else ("有代码" if s.has_code else "有图表")
                img_hints_lines.append(f"  Slide {s.index+1}: `{rel}` — {note}")
    image_hints = "\n".join(img_hints_lines) or "  (no images for this segment)"

    if has_transcript:
        # Transcript block: [MM:SS Slide N「Title」] transcript...
        transcript_lines = []
        for s in slides:
            cs = compact_by_idx.get(s.index)
            if cs and cs.get("transcript", "").strip():
                mm = int(cs["start"] // 60)
                ss = int(cs["start"] % 60)
                tx = cs["transcript"][:MAX_TRANSCRIPT_CHARS]
                transcript_lines.append(
                    f"[{mm:02d}:{ss:02d} Slide {s.index+1}「{s.label}」]\n{tx}")
        transcript_block = "\n\n".join(transcript_lines) or _P("no_transcript")

        return _P("chunk").format(
            course_name=course_name, lec_num=lec_num, lec_title=lec_title,
            slide_outline=slide_outline, transcript_block=transcript_block,
            image_hints=image_hints, chunk_idx=chunk_idx, chunk_title=chunk_title,
            detail=detail, detail_instruction=_detail_instr(detail),
        )
    else:
        return _P("slide_only").format(
            course_name=course_name, lec_num=lec_num, lec_title=lec_title,
            slide_outline=slide_outline, image_hints=image_hints,
            chunk_idx=chunk_idx, chunk_title=chunk_title,
            detail=detail, detail_instruction=_detail_instr(detail),
        )


# ── Section-by-section generation ────────────────────────────────────────────

def _section_path(sections_dir: Path, lec_num: int, ci: int, file_idx: int = 1) -> Path:
    if file_idx == 1:
        return sections_dir / f"L{lec_num:02d}_S{ci:02d}.md"
    return sections_dir / f"L{lec_num:02d}_F{file_idx:02d}_S{ci:02d}.md"


def generate_section(
    lec_num: int,
    lec_title: str,
    course_name: str,
    chunk: list[SlideInfo],
    ci: int,
    ld: "LectureData",
    out_dir: Path,
    sections_dir: Path,
    detail: int,
    has_transcript: bool,
    bar: tqdm | None = None,
    force: bool = False,
) -> str:
    """Generate (or load from cache) one section and save it to sections_dir."""
    sec_file = _section_path(sections_dir, lec_num, ci, ld.file_idx)

    if not force and sec_file.exists() and sec_file.stat().st_size > 50:
        if bar:
            fi = f"F{ld.file_idx} " if ld.file_idx > 1 else ""
            bar.set_postfix_str(f"L{lec_num}{fi}§{ci} cached")
        return sec_file.read_text(encoding="utf-8")

    chunk_title = _chunk_title(chunk)
    fi_tag = f"F{ld.file_idx} " if ld.file_idx > 1 else ""
    if bar:
        bar.set_postfix_str(f"L{lec_num}{fi_tag}§{ci}/{chunk_title[:22]} generating")
    tqdm.write(f"  → L{lec_num}{fi_tag}§{ci} [{chunk_title[:40]}]  slides {chunk[0].index+1}–{chunk[-1].index+1}  calling {NOTE_MODEL}…")

    # Render only this chunk's slide images (lazy, avoids OOM)
    img_map = ld.render_chunk_images([s.index for s in chunk])

    user = _build_chunk_prompt(
        slides=chunk,
        compact_by_idx=ld.compact_by_idx,
        img_cache=ld.img_cache,
        img_map=img_map,
        out_dir=out_dir,
        course_name=course_name,
        lec_num=lec_num,
        lec_title=lec_title,
        chunk_idx=ci,
        chunk_title=chunk_title,
        detail=detail,
        has_transcript=has_transcript,
        source=ld.source,
    )

    import time as _time
    _t0 = _time.monotonic()
    draft = _call(NOTE_MODEL, _P("system"), user, _max_tokens(detail))
    tqdm.write(f"     ✓ {len(draft):,} chars  ({_time.monotonic()-_t0:.0f}s)")

    if not draft:
        tqdm.write(f"  [warn] Empty draft for L{lec_num} §{ci} — skipping")

    if VERIFY_NOTES and draft:
        terms = set()
        for s in chunk:
            for t in re.findall(r"\b[A-Z][a-zA-Z]{3,}\b|\b[A-Z]{3,}\b", s.text):
                terms.add(t)
        term_list = ", ".join(sorted(terms)[:30])
        v_user = _P("verify").format(term_list=term_list, draft=draft[:2500])
        v_result = _call(VERIFY_MODEL, "", v_user, 1500)
        if not v_result.strip().upper().startswith("APPROVED"):
            if len(v_result) > len(draft) * 0.3:
                draft = v_result
            else:
                tqdm.write(f"  [warn] Verifier suspicious response, keeping draft")

    heading = f"### {lec_num}.{ci} {chunk_title}"
    content = f"{heading}\n\n{draft}"
    sec_file.write_text(content, encoding="utf-8")
    return content


def generate_lecture(
    lec_num: int,
    lec_title: str,
    course_name: str,
    ld: "LectureData",
    out_dir: Path,
    sections_dir: Path,
    detail: int,
    fmt: str,
    bar: tqdm | None = None,
    force: bool = False,
) -> str:
    """Generate all sections for one lecture, saving each to sections_dir."""
    has_transcript = bool(ld.compact_by_idx)
    chunks = [ld.slides[i:i+CHAPTER_SIZE]
              for i in range(0, len(ld.slides), CHAPTER_SIZE)]
    parts: list[str] = []
    for ci, chunk in enumerate(chunks, start=1):
        content = generate_section(
            lec_num=lec_num,
            lec_title=lec_title,
            course_name=course_name,
            chunk=chunk,
            ci=ci,
            ld=ld,
            out_dir=out_dir,
            sections_dir=sections_dir,
            detail=detail,
            has_transcript=has_transcript,
            bar=bar,
            force=force,
        )
        parts.append(content)
        if bar:
            bar.update(1)
    return "\n\n".join(parts)


# ── Self-scoring ──────────────────────────────────────────────────────────────

def _key_terms(text: str, n: int = 12) -> list[str]:
    tokens = re.findall(r"\b[A-Z][a-zA-Z_]{2,}\b|\b[A-Z_]{3,}\b", text)
    cs = re.findall(
        r"\b(mutex|semaphore|thread|process|deadlock|fork|kernel|scheduler|"
        r"critical.section|race.condition|barrier|starvation|spinlock|heap|"
        r"socket|packet|routing|subnet|protocol|checksum|TCP|UDP|ACK|NAK|"
        r"Viterbi|pthread|register|pipeline|cache|interrupt|syscall)\b",
        text, re.IGNORECASE,
    )
    combined = [t.lower() for t in tokens + cs]
    seen: set[str] = set()
    result: list[str] = []
    for t in sorted(set(combined), key=len, reverse=True):
        if t not in seen:
            seen.add(t); result.append(t)
        if len(result) >= n: break
    return result


def self_score(all_slides: list[SlideInfo], full_notes: str,
               all_compact: list[dict]) -> dict:
    notes_lower = full_notes.lower()

    # Coverage based on word count vs expected
    word_count      = len(full_notes.split())
    expected_words  = len(all_slides) * MIN_NOTE_WORDS_PER_SLIDE
    coverage_score  = min(word_count / expected_words, 1.0) * 10 if expected_words else 10

    # Terminology
    term_hits, term_total = 0, 0
    for s in all_slides:
        terms = _key_terms(s.text)
        if terms:
            term_total += len(terms)
            term_hits  += sum(1 for t in terms if t in notes_lower)
    terminology_score = (term_hits / term_total * 10) if term_total else 10

    # Callouts
    compact_map = {(c["slide"] - 1): c for c in all_compact}
    callouts_needed  = 0
    callouts_written = len(re.findall(r"\[!IMPORTANT\]", full_notes))
    for s in all_slides:
        cs = compact_map.get(s.index)
        if cs:
            tx = cs.get("transcript", "").lower()
            if any(w in tx for w in ("important", "remember", "exam", "key point",
                                      "critical", "note that", "this will be")):
                callouts_needed += 1
    callout_score = (min(callouts_written, callouts_needed) / callouts_needed * 10) \
                    if callouts_needed else 10

    # Code blocks
    code_needed  = sum(1 for s in all_slides if s.has_code)
    code_written = len(re.findall(r"```\w+", full_notes))
    code_score   = (min(code_written, code_needed) / code_needed * 10) if code_needed else 10

    overall = (SCORE_WEIGHTS["coverage"]    * coverage_score
             + SCORE_WEIGHTS["terminology"] * terminology_score
             + SCORE_WEIGHTS["callouts"]    * callout_score
             + SCORE_WEIGHTS["code_blocks"] * code_score)

    return {
        "coverage":    round(coverage_score,    1),
        "terminology": round(terminology_score, 1),
        "callouts":    round(callout_score,      1),
        "code_blocks": round(code_score,         1),
        "overall":     round(overall,            2),
        "stats": {
            "note_words":       word_count,
            "expected_words":   expected_words,
            "term_hits":        term_hits,
            "term_total":       term_total,
            "callouts_written": callouts_written,
            "callouts_needed":  callouts_needed,
            "code_blocks":      code_written,
            "code_slides":      code_needed,
        },
    }


def _print_score(scores: dict, label: str) -> None:
    st = scores.get("stats", {})
    tqdm.write(f"\n  ┌──────────────────────────────────────────────┐")
    tqdm.write(f"  │ 自评分: {label[:36]:36s}│")
    tqdm.write(f"  ├──────────────────────────────────────────────┤")
    tqdm.write(f"  │ 覆盖率     {scores['coverage']:4.1f}/10  "
               f"({st.get('note_words','?')} / ~{st.get('expected_words','?')} words) │")
    tqdm.write(f"  │ 术语准确率 {scores['terminology']:4.1f}/10  "
               f"({st.get('term_hits','?')}/{st.get('term_total','?')} terms)       │")
    tqdm.write(f"  │ 重点标注   {scores['callouts']:4.1f}/10  "
               f"({st.get('callouts_written','?')}/{st.get('callouts_needed','?')} callouts)     │")
    tqdm.write(f"  │ 代码块     {scores['code_blocks']:4.1f}/10  "
               f"({st.get('code_blocks','?')}/{st.get('code_slides','?')} code slides)   │")
    tqdm.write(f"  │ 综合评分   {scores['overall']:4.2f}/10                              │")
    tqdm.write(f"  └──────────────────────────────────────────────┘")


# ── LectureData ───────────────────────────────────────────────────────────────

class LectureData:
    def __init__(self, num: int, slide_path: Path, alignment_path: Path | None,
                 file_idx: int = 1, source: str = "slides",
                 frame_dir: Path | None = None):
        self.num            = num
        self.file_idx       = file_idx
        self.slide_path     = slide_path
        self.alignment_path = alignment_path
        self.source         = source       # "slides" or "screenshare"
        self.frame_dir      = frame_dir    # directory of extracted frame PNGs
        self.slides:         list[SlideInfo] = []
        self.compact:        dict            = {}
        self.compact_slides: list[dict]      = []
        self.img_cache:      dict            = {}
        self.img_map:        dict[int, Path] = {}

    def load(self, out_dir: Path) -> None:
        if self.source == "screenshare" and self.frame_dir:
            self._load_from_frames(out_dir)
        else:
            self.slides = _load_slides(self.slide_path)
        if self.alignment_path and self.alignment_path.exists():
            self.compact = alignment_parser.parse(self.alignment_path)
            self.compact_slides = self.compact.get("slides", [])
        if self.source != "screenshare":
            cache_f = self.slide_path.parent / f"{self.slide_path.name}.image_cache.json"
            if cache_f.exists():
                with open(cache_f) as f:
                    self.img_cache = json.load(f)
        # img_map is populated lazily per chunk to avoid rendering all slides at once
        self._out_dir = out_dir

    def _load_from_frames(self, out_dir: Path) -> None:
        """Load 'slides' from extracted video frames instead of PDF/PPTX."""
        frame_pngs = sorted(self.frame_dir.glob("frame_*.png"))
        self.slides = []
        for i, png in enumerate(frame_pngs):
            # For screen share frames, the text content is minimal (just the label).
            # The actual content comes from the transcript in the alignment.
            label = f"Frame {i + 1}"
            text = label  # placeholder; transcript provides the real content
            self.slides.append(SlideInfo(i, label, text))
            # Pre-populate img_map since frames are already extracted
            self.img_map[i] = png

    def render_chunk_images(self, slide_indices: list[int]) -> dict[int, Path]:
        """Render only the slides in this chunk into images/L{num}[_F{idx}]/, cache results."""
        if self.source == "screenshare":
            # For screen share, frames are already extracted — just copy/link
            # them into the notes images/ directory for consistent path references
            if self.file_idx == 1:
                img_dir = self._out_dir / "images" / f"L{self.num:02d}"
            else:
                img_dir = self._out_dir / "images" / f"L{self.num:02d}_F{self.file_idx:02d}"
            img_dir.mkdir(parents=True, exist_ok=True)
            for i in slide_indices:
                if i in self.img_map:
                    src = self.img_map[i]
                    dst = img_dir / f"slide_{i + 1:03d}.png"
                    if not dst.exists():
                        import shutil
                        shutil.copy2(str(src), str(dst))
                    self.img_map[i] = dst
            return {i: self.img_map[i] for i in slide_indices if i in self.img_map}

        if self.file_idx == 1:
            img_dir = self._out_dir / "images" / f"L{self.num:02d}"
        else:
            img_dir = self._out_dir / "images" / f"L{self.num:02d}_F{self.file_idx:02d}"
        needed = [i for i in slide_indices if i not in self.img_map]
        if needed:
            new = render_slide_images(self.slide_path, img_dir, needed)
            self.img_map.update(new)
        return {i: self.img_map[i] for i in slide_indices if i in self.img_map}

    @property
    def title(self) -> str:
        # Prefer alignment lecture title (e.g. "CS3210 e-Lecture on Processes and Threads")
        if self.compact.get("lecture"):
            t = self.compact["lecture"]
            # Strip course prefix like "CS3210 e-Lecture on "
            t = re.sub(r"^CS\d+\s+e-Lecture\s+on\s+", "", t)
            # Strip trailing speaker credit "(by ...)"
            t = re.sub(r"\s*\(by .+\)$", "", t)
            return t.strip()
        if self.source == "screenshare" and self.frame_dir:
            stem = self.frame_dir.name
            return stem.replace("-", " ").replace("_", " ")
        # Fall back to filename stem, removing "L02-" prefix
        stem = re.sub(r"^[Ll]\d+[-_\s]+", "", self.slide_path.stem)
        return stem.replace("-", " ").replace("_", " ") or self.slide_path.stem

    @property
    def compact_by_idx(self) -> dict[int, dict]:
        return {(c["slide"] - 1): c for c in self.compact_slides}


# ── Merge sections into final note ────────────────────────────────────────────

def merge_sections(
    course_name: str,
    lectures: list["LectureData"],
    sections_dir: Path,
    out_path: Path,
    all_slides: list[SlideInfo],
    all_compact: list[dict],
) -> tuple[Path, dict]:
    """Read all saved section files and merge into one final Markdown note."""
    from itertools import groupby

    note_sections: list[str] = []

    for lec_num, ld_iter in groupby(lectures, key=lambda x: x.num):
        ld_group   = list(ld_iter)
        multi_file = len(ld_group) > 1
        lec_heading = f"## Lecture {lec_num} — {ld_group[0].title}"
        lec_parts: list[str] = []

        for ld in ld_group:
            n_chunks   = max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
            file_parts: list[str] = []
            for ci in range(1, n_chunks + 1):
                sec_file = _section_path(sections_dir, ld.num, ci, ld.file_idx)
                if sec_file.exists() and sec_file.stat().st_size > 50:
                    file_parts.append(sec_file.read_text(encoding="utf-8"))
                else:
                    fi = f" F{ld.file_idx}" if multi_file else ""
                    tqdm.write(f"  [warn] Missing section L{ld.num}{fi} §{ci} — skipping")

            if file_parts:
                if multi_file:
                    part_heading = f"### Part {ld.file_idx}: {ld.slide_path.stem}"
                    lec_parts.append(f"{part_heading}\n\n" + "\n\n".join(file_parts))
                else:
                    lec_parts.extend(file_parts)

        if lec_parts:
            note_sections.append(f"{lec_heading}\n\n" + "\n\n".join(lec_parts))

    # Exam notes — generated from the merged content
    tqdm.write("  Generating exam notes…")
    summary = "\n\n---\n\n".join(note_sections)[:8000]
    exam_section_file = sections_dir / "exam_notes.md"
    if exam_section_file.exists() and exam_section_file.stat().st_size > 50:
        exam_md = exam_section_file.read_text(encoding="utf-8")
    else:
        exam_md = _call(NOTE_MODEL, _P("system"),
                        _P("exam").format(course_name=course_name, summary=summary),
                        4000)
        exam_section_file.write_text(exam_md, encoding="utf-8")
    note_sections.append(exam_md)

    now = datetime.now().strftime("%Y-%m-%dT%H:%M:%S+08:00")
    front = (f"---\ntitle: {course_name}\ndate: {now}\n"
             f"description: Lecture notes generated by auto_note\n"
             f"categories:\n    - tech\n---\n\n")
    full_notes = front + f"# {course_name} Notes\n\n" + "\n\n------\n\n".join(note_sections)

    # ── Image filter agent pass ────────────────────────────────────────────────
    tqdm.write("  Running image filter pass…")
    full_notes, _, _ = filter_images_pass(full_notes, out_path.parent, lectures)

    out_path.write_text(full_notes, encoding="utf-8")
    tqdm.write(f"\n  Merged → {out_path}  ({len(full_notes):,} chars)")

    scores = self_score(all_slides, full_notes, all_compact)
    _print_score(scores, out_path.name)
    score_path = out_path.with_suffix(".score.json")
    with open(score_path, "w") as f:
        json.dump(scores, f, indent=2)

    return out_path, scores


# ── Generate full course notes ────────────────────────────────────────────────

def generate_course_notes(
    course_name: str,
    lectures: list[LectureData],
    out_path: Path,
    detail: int = DETAIL_LEVEL,
    fmt: str = OUTPUT_FORMAT,
    force: bool = False,
) -> tuple[Path, dict]:

    out_dir = out_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    # sections/ sub-directory: each chunk saved as L{N}_S{ci}.md
    sections_dir = out_dir / "sections"
    sections_dir.mkdir(exist_ok=True)

    # Load all lectures first to get accurate chunk counts
    tqdm.write("  Loading lecture data…")
    for ld in lectures:
        ld.load(out_dir)

    total_chunks = sum(
        max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
        for ld in lectures
    )

    all_slides:  list[SlideInfo] = []
    all_compact: list[dict]      = []

    bar = tqdm(total=total_chunks, desc=f"{course_name}", unit="section",
               bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} sections [{elapsed}<{remaining}]")

    # ── Phase 1: generate each section independently ──────────────────────────
    for ld in lectures:
        all_slides.extend(ld.slides)
        all_compact.extend(ld.compact_slides)
        generate_lecture(
            lec_num=ld.num,
            lec_title=ld.title,
            course_name=course_name,
            ld=ld,
            out_dir=out_dir,
            sections_dir=sections_dir,
            detail=detail,
            fmt=fmt,
            bar=bar,
            force=force,
        )

    bar.close()

    # ── Phase 2: merge all sections into one file ─────────────────────────────
    tqdm.write("\n  Merging sections…")
    return merge_sections(
        course_name=course_name,
        lectures=lectures,
        sections_dir=sections_dir,
        out_path=out_path,
        all_slides=all_slides,
        all_compact=all_compact,
    )


# ── Per-video note generation ─────────────────────────────────────────────────

def generate_per_video_notes(
    course_name: str,
    lectures: list[LectureData],
    out_dir: Path,
    detail: int = DETAIL_LEVEL,
    fmt: str = OUTPUT_FORMAT,
    force: bool = False,
) -> list[tuple[Path, dict]]:
    """Generate one note file per lecture/video instead of a single merged file.

    Each video gets its own Markdown file: <course>_L<N>_notes.md
    Returns a list of (path, scores) tuples for each generated note.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    sections_dir = out_dir / "sections"
    sections_dir.mkdir(exist_ok=True)

    tqdm.write("  Loading lecture data…")
    for ld in lectures:
        ld.load(out_dir)

    results: list[tuple[Path, dict]] = []

    for ld in lectures:
        lec_num = ld.num
        lec_title = ld.title
        ext = ".mdx" if fmt == "mdx" else ".md"
        note_path = out_dir / f"{course_name}_L{lec_num:02d}_notes{ext}"

        tqdm.write(f"\n  ═══ Lecture {lec_num}: {lec_title} ═══")

        n_chunks = max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
        bar = tqdm(total=n_chunks, desc=f"L{lec_num}", unit="section",
                   bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]")

        has_transcript = bool(ld.compact_by_idx)
        chunks = [ld.slides[i:i + CHAPTER_SIZE]
                  for i in range(0, len(ld.slides), CHAPTER_SIZE)]

        parts: list[str] = []
        for ci, chunk in enumerate(chunks, start=1):
            content = generate_section(
                lec_num=lec_num,
                lec_title=lec_title,
                course_name=course_name,
                chunk=chunk,
                ci=ci,
                ld=ld,
                out_dir=out_dir,
                sections_dir=sections_dir,
                detail=detail,
                has_transcript=has_transcript,
                bar=bar,
                force=force,
            )
            parts.append(content)
            bar.update(1)
        bar.close()

        # Build per-video note
        from datetime import datetime
        now = datetime.now().strftime("%Y-%m-%dT%H:%M:%S+08:00")
        front = (f"---\ntitle: {course_name} — Lecture {lec_num}\ndate: {now}\n"
                 f"description: Lecture notes generated by auto_note\n"
                 f"categories:\n    - tech\n---\n\n")
        heading = f"# {course_name} — Lecture {lec_num}: {lec_title}\n\n"
        full_notes = front + heading + "\n\n".join(parts)

        # Image filter
        tqdm.write("  Running image filter pass…")
        full_notes, _, _ = filter_images_pass(full_notes, out_dir, [ld])

        note_path.write_text(full_notes, encoding="utf-8")
        tqdm.write(f"  Saved → {note_path}  ({len(full_notes):,} chars)")

        scores = self_score(ld.slides, full_notes, ld.compact_slides)
        _print_score(scores, note_path.name)
        score_path = note_path.with_suffix(".score.json")
        with open(score_path, "w") as f:
            json.dump(scores, f, indent=2)

        results.append((note_path, scores))

    return results


# ── Iteration ─────────────────────────────────────────────────────────────────

def generate_with_iteration(
    course_name: str,
    lectures: list[LectureData],
    out_path: Path,
    fmt: str = OUTPUT_FORMAT,
    max_rounds: int = 3,
) -> Path:
    detail     = DETAIL_LEVEL
    best_path  = out_path
    best_score = 0.0

    for rnd in range(1, max_rounds + 1):
        print(f"\n{'#'*60}")
        print(f"# ROUND {rnd}/{max_rounds}   detail={detail}   target={QUALITY_TARGET}")
        print(f"{'#'*60}")

        path, scores = generate_course_notes(
            course_name, lectures, out_path,
            detail=detail, fmt=fmt,
            force=(rnd > 1),   # re-generate sections on subsequent rounds
        )
        overall = scores["overall"]
        if overall > best_score:
            best_score = overall
            best_path  = path

        if overall >= QUALITY_TARGET:
            tqdm.write(f"\n  ✓ 达到目标 {QUALITY_TARGET} (得分={overall:.2f})")
            break
        if rnd < max_rounds:
            tqdm.write(f"\n  得分 {overall:.2f} < {QUALITY_TARGET}，提升详细度后重试…")
            versioned = path.with_name(f"{path.stem}_r{rnd}{path.suffix}")
            path.rename(versioned)
            detail = min(detail + 2, 10)

    tqdm.write(f"\n  最佳: {best_path}  (得分={best_score:.2f})")
    return best_path


# ── Auto-discovery ────────────────────────────────────────────────────────────

def _find_alignment(slide_path: Path, course_dir: Path) -> Path | None:
    align_dir = course_dir / "alignment"
    if not align_dir.exists():
        return None
    stem = slide_path.stem.lower()
    best_sc, best = 0.0, None
    for f in align_dir.glob("*.json"):
        if f.name.endswith(".compact.json"):
            continue
        ta = set(stem.replace("-", " ").split())
        tb = set(f.stem.lower().replace("-", " ").split())
        if not ta or not tb:
            continue
        sc = len(ta & tb) / len(ta | tb)
        if sc > best_sc:
            best_sc, best = sc, f
    return best if best_sc > 0.05 else None


def _discover_screenshare_lectures(course_dir: Path) -> list[LectureData]:
    """Discover lectures from screen share frame directories.

    Checks the manifest to find screen share (SS) videos, then looks for
    extracted frames in <course_dir>/frames/<video_stem>/.
    """
    frames_dir = course_dir / "frames"
    if not frames_dir.exists():
        return []

    lectures: list[LectureData] = []
    frame_subdirs = sorted([d for d in frames_dir.iterdir() if d.is_dir()])
    for i, fdir in enumerate(frame_subdirs, start=1):
        frame_pngs = sorted(fdir.glob("frame_*.png"))
        if not frame_pngs:
            continue
        # Check for alignment
        align_path = course_dir / "alignment" / f"{fdir.name}.json"
        align = align_path if align_path.exists() else None
        # Verify alignment source is screenshare
        if align:
            try:
                with open(align, encoding="utf-8") as f:
                    data = json.load(f)
                if data.get("source") != "screenshare":
                    continue  # This alignment is from traditional slide matching
            except Exception:
                pass
        else:
            continue  # No alignment yet — need frame_extractor to run first

        lectures.append(LectureData(
            num=i,
            slide_path=fdir,   # frame_dir serves as slide_path for compatibility
            alignment_path=align,
            file_idx=1,
            source="screenshare",
            frame_dir=fdir,
        ))
    return lectures


def _discover_lectures(course_dir: Path) -> list[LectureData]:
    exts    = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}

    # ── Check for screen share lectures first ────────────────────────────────
    ss_lectures = _discover_screenshare_lectures(course_dir)
    if ss_lectures:
        tqdm.write(f"  Found {len(ss_lectures)} screen share lecture(s) from extracted frames")
        return ss_lectures

    # ── Traditional slide-based discovery ─────────────────────────────────────

    # Common lecture-slide subfolder names, checked in priority order.
    _LECTURE_SUBDIRS = [
        "LectureNotes", "Lecture Slides", "Lecture Notes",
        "Lectures", "Slides", "lecture_notes", "lecture_slides",
    ]

    mat_dir = course_dir / "materials"
    if not mat_dir.exists():
        print(f"[error] Materials directory not found: {mat_dir}")
        print("  Run 'Download materials' first before generating notes.")
        sys.exit(1)
    lecture_subdir: Path | None = None
    for name in _LECTURE_SUBDIRS:
        candidate = mat_dir / name
        if candidate.exists():
            lecture_subdir = candidate
            break

    if lecture_subdir is not None:
        slide_files = sorted([
            p for p in lecture_subdir.rglob("*")
            if p.suffix.lower() in exts and "image_cache" not in p.name
        ])
    else:
        # Recursively search all of materials/ for lecture-like files.
        # Include a file if it has a lecture number in its name OR lives
        # in a folder whose name suggests lectures.
        _LECTURE_FOLDER_PAT = re.compile(
            r"lecture|slides|notes", re.IGNORECASE)
        _LECTURE_FILE_PAT = re.compile(
            r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_ ]?0*\d+", re.IGNORECASE)

        all_candidates = sorted([
            p for p in mat_dir.rglob("*")
            if p.is_file() and p.suffix.lower() in exts
            and "image_cache" not in p.name
        ])

        slide_files = []
        for p in all_candidates:
            # Always include files at the root of materials/
            if p.parent == mat_dir:
                slide_files.append(p)
                continue
            # Include if filename looks like a lecture
            if _LECTURE_FILE_PAT.search(p.stem):
                slide_files.append(p)
                continue
            # Include if any parent folder name suggests lectures
            rel_parts = p.relative_to(mat_dir).parts[:-1]  # folder names
            if any(_LECTURE_FOLDER_PAT.search(part) for part in rel_parts):
                slide_files.append(p)
                continue

    # Group files by lecture number; files without a number get a unique high index
    from collections import defaultdict
    by_num: dict[int, list[Path]] = defaultdict(list)
    auto_num = 1000
    for sp in slide_files:
        m = re.search(r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_ ]?0*(\d+)", sp.stem)
        if m:
            lec_num = int(m.group(1))
        else:
            lec_num  = auto_num
            auto_num += 1
        by_num[lec_num].append(sp)

    lectures: list[LectureData] = []
    for lec_num in sorted(by_num.keys()):
        files = sorted(by_num[lec_num])   # alphabetical within same lecture number
        for file_idx, sp in enumerate(files, start=1):
            align = _find_alignment(sp, course_dir)
            lectures.append(LectureData(lec_num, sp, align, file_idx=file_idx))
    return lectures


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Generate course lecture notes")
    parser.add_argument("--course",      metavar="ID")
    parser.add_argument("--slides",      metavar="PATH")
    parser.add_argument("--alignment",   metavar="PATH")
    parser.add_argument("--lecture-num", metavar="N", type=int, default=1)
    parser.add_argument("--course-name", metavar="NAME", default="")
    parser.add_argument("--out",         metavar="PATH")
    parser.add_argument("--detail",      metavar="0-10", type=int, default=DETAIL_LEVEL)
    parser.add_argument("--format",      metavar="md|mdx", default=OUTPUT_FORMAT,
                        choices=["md", "mdx"])
    parser.add_argument("--iterate",     action="store_true")
    parser.add_argument("--merge-only",  action="store_true",
                        help="Skip generation; just merge existing sections into final note")
    parser.add_argument("--force",       action="store_true",
                        help="Re-generate sections even if they already exist")
    parser.add_argument("--lectures",    metavar="N-N or N,N,N", default="",
                        help="Filter lectures, e.g. '1-5' or '1,2,3'")
    parser.add_argument("--per-video",  action="store_true",
                        help="Generate one note file per video/lecture instead of one combined note")
    args = parser.parse_args()

    if args.course:
        course_dir  = COURSE_DATA_DIR / args.course
        course_name = args.course_name or f"CS{args.course}"
        lectures    = _discover_lectures(course_dir)
        if args.lectures:
            sel: set[int] = set()
            for part in args.lectures.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    sel.update(range(int(a), int(b) + 1))
                elif part.isdigit():
                    sel.add(int(part))
            lectures = [l for l in lectures if l.num in sel]
        if not lectures:
            print(f"No slides found under {course_dir}"); sys.exit(1)

        print(f"Found {len(lectures)} lectures:")
        for ld in lectures:
            src = f" [{ld.source}]" if ld.source != "slides" else ""
            a  = "+ alignment" if ld.alignment_path else "(slide-only)"
            fi = f" [part {ld.file_idx}]" if ld.file_idx > 1 else ""
            print(f"  L{ld.num}{fi}{src}: {ld.slide_path.name}  {a}")

        ext_out      = ".mdx" if args.format == "mdx" else ".md"
        out_path     = Path(args.out) if args.out else \
                       course_dir / "notes" / f"{course_name}_notes{ext_out}"
        sections_dir = out_path.parent / "sections"

        if args.per_video:
            # Per-video mode: one note per lecture/video
            out_dir = out_path.parent
            generate_per_video_notes(
                course_name, lectures, out_dir,
                detail=args.detail, fmt=args.format, force=args.force,
            )
            return

        if args.merge_only:
            # Load lectures so titles are available
            out_path.parent.mkdir(parents=True, exist_ok=True)
            for ld in lectures:
                ld.load(out_path.parent)
            all_slides  = [s for ld in lectures for s in ld.slides]
            all_compact = [c for ld in lectures for c in ld.compact_slides]
            merge_sections(course_name, lectures, sections_dir, out_path,
                           all_slides, all_compact)
            return

        if args.iterate:
            generate_with_iteration(course_name, lectures, out_path, fmt=args.format)
        else:
            generate_course_notes(course_name, lectures, out_path,
                                  detail=args.detail, fmt=args.format,
                                  force=args.force)
        return

    if not args.slides:
        parser.error("Provide --course or --slides")

    sp = Path(args.slides)
    if not sp.exists():
        print(f"[error] Not found: {sp}"); sys.exit(1)

    align    = Path(args.alignment) if args.alignment else None
    ld       = LectureData(args.lecture_num, sp, align)
    name     = args.course_name or sp.parent.parent.parent.name
    ext_out  = ".mdx" if args.format == "mdx" else ".md"
    out_path = Path(args.out) if args.out else \
               sp.parent.parent.parent / "notes" / f"{name}_notes{ext_out}"

    fn = generate_with_iteration if args.iterate else \
         lambda *a, **kw: generate_course_notes(*a, **kw)
    fn(name, [ld], out_path, fmt=args.format,
       **({} if args.iterate else {"detail": args.detail}))


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n[info] Interrupted by user.")
        sys.exit(0)
    except SystemExit:
        raise
    except Exception as _exc:
        import traceback
        print(f"\n[error] Unexpected error: {_exc}")
        traceback.print_exc()
        sys.exit(1)
